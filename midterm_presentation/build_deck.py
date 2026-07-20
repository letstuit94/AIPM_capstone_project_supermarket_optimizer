#!/usr/bin/env python3
"""NutriWise midterm deck — corporate design (sage green, sans-serif).
Charts + .pptx, 13 slides (+1 backup) incl. live demo.

v2: compressed for a 10-min slot incl. demo; integrates the ML Step-2
gold-set results (Recall@5 on real receipts) and adds a data-maturity
("honest by design") slide. EDA numbers still come from the repo
(bls_off_judgments.json); the gold-set numbers are the reported Step-2
figures from the ML notebooks (goldset_extend / goldset_eval / off_fallback)."""
import json, collections, os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------- corporate palette
# sourced from frontend/src/index.css : one accent (sage), neutral ink/canvas
SAGE     = "7C9A6A"   # --accent
SAGE_DK  = "5F7B4F"   # darker sage for emphasis
SAGE_LT  = "B7C9A9"   # partial / soft mark
SAGE_SOFT= "EEF2EA"   # --accent-soft (panels)
CLAY     = "B36A4A"   # the single restrained "negative" warm tone
SAND     = "CFC9BA"   # neutral no-match (warm gray)
INK      = "1D1D21"   # --ink
INK_SOFT = "6E6F74"   # --ink-soft
CANVAS   = "FAFAFA"   # --canvas
SURFACE  = "FFFFFF"   # --surface
LINE     = "E5E4E0"   # hairline border
WHITE    = "FFFFFF"
FONT     = "Arial"    # system sans-serif; matches app's --font-sans (no serif)

def C(h): return RGBColor.from_string(h)

SCRATCH = os.path.dirname(os.path.abspath(__file__))
ROOT    = "/Users/jenniferrake/Desktop/GitHub_Bootcamp/AIPM_capstone_project_supermarket_optimizer"
OUTDIR  = os.path.join(ROOT, "midterm_presentation")
CH = lambda n: os.path.join(SCRATCH, n)

TOTAL_PAGES = 13   # main slides (title..thanks); backup slide is unnumbered

# ---------------------------------------------------- data (EDA + old baseline)
d = json.load(open(os.path.join(ROOT, "bls_off_judgments.json")))
N = len(d)
off = collections.Counter(x["off_verdict"] for x in d)
bls = collections.Counter(x["bls_verdict"] for x in d)

GROUP = {
    "Gemüse":"Fruit & Veg","Obst":"Fruit & Veg","Obst & Gemüse":"Fruit & Veg",
    "Milchprodukte":"Dairy, Cheese & Eggs","Käse":"Dairy, Cheese & Eggs",
    "Milchprodukte & Alternativen":"Dairy, Cheese & Eggs","Eier":"Dairy, Cheese & Eggs",
    "Backwaren":"Bread & Grains","Müsli & Cerealien":"Bread & Grains",
    "Nudeln & Reis":"Bread & Grains","grain":"Bread & Grains",
    "Fleisch":"Meat & Fish","Fleisch/Wurst":"Meat & Fish","Fisch":"Meat & Fish",
    "Hülsenfrüchte":"Legumes & Protein","protein":"Legumes & Protein",
    "Getränke":"Drinks & Tea","Tee":"Drinks & Tea",
}
grp = collections.Counter()
for x in d:
    grp[GROUP.get(x["category"], "Other / processed")] += 1

# ---------------------------------------------------- matplotlib defaults
plt.rcParams.update({
    "font.family":"sans-serif","font.sans-serif":["Arial","Helvetica","DejaVu Sans"],
    "font.size":15,"text.color":"#"+INK,"axes.edgecolor":"#"+LINE,
    "axes.labelcolor":"#"+INK,"xtick.color":"#"+INK_SOFT,"ytick.color":"#"+INK,
    "figure.facecolor":"#"+SURFACE,"axes.facecolor":"#"+SURFACE,"savefig.facecolor":"#"+SURFACE,
})
def strip(ax, keep_x=False):
    for s in ["top","right","left"]: ax.spines[s].set_visible(False)
    if not keep_x:
        ax.spines["bottom"].set_visible(False); ax.tick_params(bottom=False)
    ax.tick_params(left=False)

# ---- CHART 1: EDA food-group distribution
labels=[k for k,_ in grp.most_common()]; vals=[v for _,v in grp.most_common()]
fig,ax=plt.subplots(figsize=(6.6,4.3),dpi=200)
y=list(range(len(labels)))[::-1]
bars=ax.barh(y,vals,color="#"+SAGE,height=0.66); bars[0].set_color("#"+SAGE_DK)
ax.set_yticks(y); ax.set_yticklabels(labels,fontsize=13); ax.set_xticks([])
for yi,v in zip(y,vals):
    ax.text(v+0.3,yi,str(v),va="center",ha="left",fontsize=13,color="#"+INK,fontweight="bold")
strip(ax); ax.set_xlim(0,max(vals)+3)
ax.set_title("What people actually buy — 76 unique items",
             fontsize=14.5,fontweight="bold",color="#"+INK,loc="left",pad=12)
plt.tight_layout(); plt.savefig(CH("eda_categories.png"),bbox_inches="tight"); plt.close()

# ---- CHART 2: EDA long-tail
occ=collections.Counter(x["occurrences"] for x in d)
buckets={"once":occ[1],"2×":occ[2],"3–4×":occ[3]+occ[4],
         "5–9×":sum(occ[k] for k in range(5,10)),
         "10×+":sum(v for k,v in occ.items() if k>=10)}
fig,ax=plt.subplots(figsize=(6.6,4.3),dpi=200)
bl=list(buckets); bv=list(buckets.values())
bars=ax.bar(bl,bv,color="#"+SAGE_LT,width=0.62); bars[0].set_color("#"+SAGE)
for b,v in zip(bars,bv):
    ax.text(b.get_x()+b.get_width()/2,v+0.4,str(v),ha="center",va="bottom",
            fontsize=13,fontweight="bold",color="#"+INK)
strip(ax,keep_x=True); ax.set_yticks([]); ax.set_ylim(0,max(bv)+4)
ax.tick_params(axis="x",labelsize=13)
ax.set_title("Real baskets are long-tail: half of all items\nappear only once",
             fontsize=14.5,fontweight="bold",color="#"+INK,loc="left",pad=12)
plt.tight_layout(); plt.savefig(CH("eda_tail.png"),bbox_inches="tight"); plt.close()

# ---- CHART 3: cleanup effect (ML Step-2, gold set of 180 real receipt lines)
# reported figures: ExpTop-1 18%→25%, Recall@5 32%→47% (raw vs cleaned text)
fig,ax=plt.subplots(figsize=(6.9,4.3),dpi=200)
groups=["Top-1 auto-pick","Right food in top 5"]
raw=[18,32]; clean=[25,47]
xs=range(len(groups)); w=0.36
b1=ax.bar([i-w/2 for i in xs],raw,width=w,color="#"+SAND,label="raw receipt text")
b2=ax.bar([i+w/2 for i in xs],clean,width=w,color="#"+SAGE_DK,label="cleaned text")
for bars in (b1,b2):
    for b in bars:
        ax.text(b.get_x()+b.get_width()/2,b.get_height()+0.8,f"{int(b.get_height())}%",
                ha="center",va="bottom",fontsize=14,fontweight="bold",color="#"+INK)
ax.set_xticks(list(xs)); ax.set_xticklabels(groups,fontsize=13)
ax.set_ylim(0,58); ax.set_yticks([]); strip(ax,keep_x=True)
ax.legend(frameon=False,fontsize=12,loc="upper left",bbox_to_anchor=(0,1.02))
ax.set_title("Cleaning the receipt text first lifts every metric",
             fontsize=14.5,fontweight="bold",color="#"+INK,loc="left",pad=14)
plt.tight_layout(); plt.savefig(CH("cleanup.png"),bbox_inches="tight"); plt.close()

# ---- CHART 4 (BACKUP): baseline — OFF vs BLS head-to-head (pilot, LLM-judged)
order=["correct","partially_correct","no_match","incorrect"]
oleg={"correct":"correct","partially_correct":"partial","no_match":"no match","incorrect":"wrong"}
colors={"correct":SAGE,"partially_correct":SAGE_LT,"no_match":SAND,"incorrect":CLAY}
rows=[("OpenFoodFacts\nidentity match",off),("BLS generic\ntable match",bls)]
fig,ax=plt.subplots(figsize=(9.2,3.5),dpi=200); ypos=[1,0]
for (name,cnt),yp in zip(rows,ypos):
    left=0
    for k in order:
        wv=cnt.get(k,0)/N*100
        ax.barh(yp,wv,left=left,color="#"+colors[k],height=0.5,edgecolor="#"+SURFACE,linewidth=1.5)
        if wv>=7:
            ax.text(left+wv/2,yp,f"{round(wv)}%",ha="center",va="center",fontsize=13,
                    fontweight="bold",color="#"+WHITE if k in("correct","incorrect") else "#"+INK)
        left+=wv
ax.set_yticks(ypos); ax.set_yticklabels([r[0] for r in rows],fontsize=13,fontweight="bold")
ax.set_xticks([]); ax.set_xlim(0,100); ax.set_ylim(-0.6,1.6); strip(ax)
handles=[plt.Rectangle((0,0),1,1,color="#"+colors[k]) for k in order]
ax.legend(handles,[oleg[k] for k in order],ncol=4,loc="upper center",
          bbox_to_anchor=(0.5,-0.08),frameon=False,fontsize=12,handlelength=1.1)
plt.tight_layout(); plt.savefig(CH("baseline.png"),bbox_inches="tight"); plt.close()
print("charts done")

# ==================================================== PPTX
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def bg(s,h): s.background.fill.solid(); s.background.fill.fore_color.rgb=C(h)

def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,shadow=False,adj=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=C(line); sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if adj is not None:
        try: sp.adjustments[0]=adj
        except Exception: pass
    if shadow:
        from pptx.oxml.ns import qn
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{})
        sh=el.makeelement(qn('a:outerShdw'),{'blurRad':'80000','dist':'30000','dir':'5400000','rotWithShape':'0'})
        clr=el.makeelement(qn('a:srgbClr'),{'val':'1D1D21'}); alp=el.makeelement(qn('a:alpha'),{'val':'12000'})
        clr.append(alp); sh.append(clr); ef.append(sh); el.append(ef)
    return sp

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space_after=4,ls=1.0,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,col,b,it) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it
            r.font.color.rgb=C(col); r.font.name=FONT
    return tb
def R(t,sz,col=INK,b=False,it=False): return (t,sz,col,b,it)

def notes(s,speaker,body): s.notes_slide.notes_text_frame.text=f"[{speaker}]\n{body}"
def kicker(s,x,y,t): txt(s,x,y,8,0.3,[[R(t.upper(),12.5,SAGE_DK,True,False)]])
def heading(s,x,y,w,t,size=29): txt(s,x,y,w,1.0,[[R(t,size,INK,True,False)]],ls=1.0)
def rule(s,x,y,w=1.5): rect(s,x,y,w,0.045,fill=SAGE)
def footer(s,page):
    txt(s,0.6,7.03,5,0.3,[[R("NutriWise",9,INK_SOFT,True,False),R("  ·  V2",9,INK_SOFT,False,False)]])
    txt(s,10.8,7.03,1.9,0.3,[[R(f"{page:02d} / {TOTAL_PAGES}",9,INK_SOFT,False,False)]],align=PP_ALIGN.RIGHT)
def logo(s,x,y,dark=False):
    rect(s,x,y,0.34,0.34,fill=SAGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.3)
    txt(s,x+0.46,y-0.03,4,0.42,[[R("NUTRIWISE",13,WHITE if dark else INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
def contentbg(s): bg(s,CANVAS)

def stat(s,x,y,w,h,big,label,sub=None,accent=SAGE):
    rect(s,x,y,w,h,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.06)
    rect(s,x,y+0.14,0.08,h-0.28,fill=accent)
    body=[[R(big,29,accent,True,False)],[R(label,12.5,INK,True,False)]]
    if sub: body.append([R(sub,10.5,INK_SOFT,False,False)])
    txt(s,x+0.3,y+0.14,w-0.42,h-0.26,body,anchor=MSO_ANCHOR.MIDDLE,space_after=3,ls=1.02)

def pill(s,x,y,w,t,fill=SAGE,col=WHITE):
    rect(s,x,y,w,0.44,fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    txt(s,x,y+0.02,w,0.4,[[R(t,11,col,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ================= SLIDE 1 — Title
s=slide(); bg(s,CANVAS)
rect(s,0,0,13.333,0.14,fill=SAGE)
rect(s,8.9,0.14,4.433,7.36,fill=SAGE_SOFT)     # soft-sage side band
logo(s,0.9,0.85)
txt(s,0.9,2.15,8.0,2.6,
    [[R("Upload a grocery receipt.",42,INK,True,False)],
     [R("Get one smarter thing",42,SAGE_DK,True,False)],
     [R("to buy next time.",42,SAGE_DK,True,False)]],ls=1.05)
rule(s,0.92,4.55,3.4)
txt(s,0.92,4.8,7.6,0.7,[[R("A receipt-first nutrition assistant — your pantry works for you: tap what you eat, don't type it.",15,INK_SOFT,False,True)]])
txt(s,0.9,5.75,7.8,1.2,
    [[R("Jennifer Rake",15,INK,True,False),R("  · Speaker 2      ",13,SAGE_DK,True,False),
      R("Stuart Kasemeier",15,INK,True,False),R("  · Speaker 1",13,INK_SOFT,True,False)],
     [R("AI Product Management Capstone  ·  Midterm Presentation  ·  Version 2",12.5,INK_SOFT,False,False)]],space_after=8)
# side band motif
for i,t in enumerate(["Onboard","Pantry","Check off","Insights"]):
    yy=2.35+i*0.92
    rect(s,9.5,yy,0.26,0.26,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,9.95,yy-0.06,2.9,0.4,[[R(t,14,INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if i<3: rect(s,9.62,yy+0.3,0.02,0.6,fill=SAGE_LT)
notes(s,"Both — Speaker 1 opens",
      "One line to open: 'Take a photo of your shopping receipt, and we tell you one smart thing to buy "
      "next time.' Say the two names. Tell them the plan: about 10 minutes, kept simple, with a quick live "
      "demo in the middle. Speaker 1 does the product half; Speaker 2 does the data half after the demo.")

# ================= SLIDE 2 — Problem
s=slide(); contentbg(s); kicker(s,0.9,0.6,"The problem"); heading(s,0.9,0.92,11,"Eating well shouldn't require a spreadsheet"); rule(s,0.92,1.62)
txt(s,0.9,2.0,6.1,3.6,
    [[R("Health-conscious adults (25–45) want to eat better — but every tool asks them to ",15.5,INK,False,False),
      R("log every meal by hand.",15.5,INK,True,False)],
     [R("",7,INK,False,False)],
     [R("That habit dies within weeks. People give up — or they guess.",15.5,INK,False,False)],
     [R("",7,INK,False,False)],
     [R("Meanwhile the one honest record of what they actually eat already exists — ",15.5,INK,False,False),
      R("the grocery receipt.",15.5,SAGE_DK,True,False)]],ls=1.14)
px=7.4; rect(s,px,1.95,5.05,3.85,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05)
txt(s,px+0.35,2.2,4.4,0.4,[[R("WHY IT MATTERS",12,SAGE_DK,True,False)]])
for i,t in enumerate(["Manual logging is the #1 reason nutrition apps get abandoned",
                      "Receipts are passive — zero extra effort for the user",
                      "One clear “next cart” decision beats another dashboard to ignore"]):
    yy=2.68+i*1.0
    rect(s,px+0.35,yy+0.03,0.34,0.34,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,px+0.35,yy-0.01,0.34,0.34,[[R(str(i+1),13,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,px+0.85,yy-0.05,3.9,0.9,[[R(t,13.5,INK,False,False)]],ls=1.05)
footer(s,2)
notes(s,"Speaker 1",
      "Simple version: people want to eat healthier, but every app makes them type in every meal — that's "
      "boring, so they quit. Our trick: the shopping receipt already shows what they bought, and they "
      "didn't have to do anything extra. That's our way in. Target user: health-minded people aged 25–45.")

# ================= SLIDE 3 — Solution & product (merged loop + real UI)
s=slide(); contentbg(s); kicker(s,0.9,0.6,"The solution"); heading(s,0.9,0.92,11.7,"One receipt in → a pantry that works for you",size=27); rule(s,0.92,1.62)
txt(s,0.9,1.62,11.5,0.4,[[R("Upload  ›  Your Pantry  ›  Check off what you ate  ›  Weekly gaps + 2 picks",13,SAGE_DK,True,False)]])
pw=2.55; ph=3.4; py=2.25; xs=[0.95,3.92,6.89,9.86]
caps=[("1 · Onboarding chat","conversational profile"),("2 · Your Pantry","home · expiry-aware"),
      ("3 · Daily log","1 tap = 1 portion"),("4 · Insights","weekly gaps + picks")]
def mpshell(x,title):
    rect(s,x,py,pw,ph,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.07)
    rect(s,x,py,pw,0.5,fill=SAGE,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.12)
    txt(s,x+0.2,py+0.05,pw-0.3,0.42,[[R(title,10,WHITE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
for k in range(3):
    txt(s,xs[k]+pw+0.02,py+1.4,0.4,0.5,[[R("›",22,SAGE,True,False)]],align=PP_ALIGN.CENTER)
for k,(t,sub) in enumerate(caps):
    txt(s,xs[k],py+ph+0.12,pw,0.55,[[R(t,11.5,INK,True,False)],[R(sub,10,INK_SOFT,False,False)]],align=PP_ALIGN.CENTER,space_after=1)

# phone 1: chat
x=xs[0]; mpshell(x,"NutriWise")
rect(s,x+0.2,py+0.72,1.55,0.4,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.45)
txt(s,x+0.33,py+0.74,1.4,0.36,[[R("Dein Ziel?",8,INK,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
rect(s,x+0.8,py+1.24,1.55,0.4,fill=SAGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.45)
txt(s,x+0.93,py+1.26,1.4,0.36,[[R("Gesünder essen",7.5,WHITE,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
for j,lab in enumerate(["Wenig","Mittel","Hoch"]):
    cxp=x+0.2+j*0.72
    rect(s,cxp,py+1.92,0.66,0.34,fill=SURFACE,line=SAGE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    txt(s,cxp,py+1.94,0.66,0.3,[[R(lab,7,SAGE_DK,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
for j in range(4):
    rect(s,x+0.2+j*0.2,py+2.6,0.14,0.05,fill=SAGE if j==0 else LINE)

# phone 2: pantry
x=xs[1]; mpshell(x,"Dein Vorrat")
rect(s,x+0.2,py+0.68,2.15,0.55,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.2)
txt(s,x+0.33,py+0.72,1.95,0.5,[[R("EMPFOHLEN",7,SAGE_DK,True,False)],[R("🥛 Trink deine Buttermilch",8.5,INK,True,False)]],space_after=1)
for j,(nm,bd,cc) in enumerate([("🍗 Hähnchenbrust","1 Tag",CLAY),("🥛 Buttermilch","2 Tage",CLAY),("🥦 Brokkoli","frisch",SAGE_DK)]):
    yb=py+1.42+j*0.6
    rect(s,x+0.2,yb,2.15,0.5,fill=SURFACE,line=LINE,lw=0.6,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.2)
    txt(s,x+0.33,yb,1.5,0.5,[[R(nm,8,INK,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+1.7,yb,0.6,0.5,[[R(bd,7,cc,True,False)]],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.RIGHT)

# phone 3: log
x=xs[2]; mpshell(x,"Tages-Log")
for j,(nm,val,frac,cc) in enumerate([("Protein","62 / 120 g",0.52,CLAY),("Ballaststoffe","11 / 30 g",0.37,CLAY),("Kalorien","1.340 / 2.100",0.64,SAGE)]):
    yb=py+0.82+j*0.72
    txt(s,x+0.2,yb,2.15,0.26,[[R(nm,7.5,INK,True,False),R("   "+val,7,INK_SOFT,False,False)]])
    rect(s,x+0.2,yb+0.3,2.15,0.14,fill=LINE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    rect(s,x+0.2,yb+0.3,2.15*frac,0.14,fill=cc,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)

# phone 4: insights
x=xs[3]; mpshell(x,"Insights")
txt(s,x+0.2,py+0.62,2.15,0.66,[[R("72",30,SAGE_DK,True,False)]],align=PP_ALIGN.CENTER)
txt(s,x+0.2,py+1.28,2.15,0.3,[[R("Wochen-Score · Gut",8,INK_SOFT,False,False)]],align=PP_ALIGN.CENTER)
for j,(nm,st,cc) in enumerate([("Protein","niedrig",CLAY),("Ballastst.","niedrig",CLAY),("Zucker","ok",SAGE_DK),("Kalorien","ok",SAGE_DK)]):
    cxp=x+0.2+(j%2)*1.1; cyp=py+1.68+(j//2)*0.6
    rect(s,cxp,cyp,1.02,0.5,fill=SURFACE,line=LINE,lw=0.6,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.2)
    txt(s,cxp+0.11,cyp,0.85,0.5,[[R(nm,7,INK,False,False)],[R(st,7.5,cc,True,False)]],space_after=0,anchor=MSO_ANCHOR.MIDDLE)

txt(s,0.9,6.4,11.5,0.4,[[R("Fully clickable prototype — light & dark.  ",11,INK_SOFT,False,True),
    R("Estimates from your groceries + taps — not medical advice.",11,SAGE_DK,True,True)]],align=PP_ALIGN.CENTER)
footer(s,3)
notes(s,"Speaker 1",
      "The whole loop in one picture. First you chat a little so we learn your goal. Then a photo of your "
      "receipt fills your 'pantry' — the list of food in your house. Every day you tap what you ate (one "
      "tap, no typing). At the end of the week we show what's missing and two easy fixes. Say plainly: "
      "these are friendly estimates, not a doctor's advice. A real clickable app exists if they want to see it.")

# ================= SLIDE 4 — Honest by design (data-maturity tiers)
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Honest by design"); heading(s,0.9,0.92,11.7,"We show only what your data actually supports",size=27); rule(s,0.92,1.62)
txt(s,0.9,1.66,11.5,0.4,[[R("The more you track, the more we can say — and we never pretend to know more than we do.",13,INK_SOFT,False,True)]])
ladder=[("DAY 0","Your targets","What to aim for, from your profile. No eating data yet.",SAGE_LT),
        ("1ST RECEIPT","Your basket","The mix of what you bought — labeled “basket,” never “diet.”",SAGE_LT),
        ("FEW TAPS","Early signs","Soft hints + a “3 / 7 days tracked” bar. No hard score yet.",SAGE),
        ("1 WEEK","Weekly score","Your main number: gaps per nutrient, allowing for meals eaten out.",SAGE_DK),
        ("30–90 DAYS","Trends","Are you improving over time — a line, not a single badge.",SAGE_DK)]
cw=2.14; gp=0.19; x0=0.9; yy=2.35; ch=2.75
for i,(tag,title,body,col) in enumerate(ladder):
    x=x0+i*(cw+gp)
    rect(s,x,yy,cw,ch,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    rect(s,x,yy,cw,0.66,fill=col,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.09)
    tc=INK if col==SAGE_LT else WHITE
    txt(s,x,yy+0.1,cw,0.46,[[R(tag,10.5,tc,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.22,yy+0.82,cw-0.4,0.5,[[R(title,15,INK,True,False)]])
    txt(s,x+0.22,yy+1.35,cw-0.4,1.3,[[R(body,11.5,INK_SOFT,False,False)]],ls=1.08)
    if i<4: txt(s,x+cw-0.02,yy+0.95,0.3,0.5,[[R("›",22,SAGE,True,False)]])
rect(s,0.9,5.5,11.53,1.05,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,1.2,5.62,11,0.85,
    [[R("Never a confident wrong number.  ",13.5,SAGE_DK,True,False),
      R("A coverage floor holds the score back until enough days are tracked, and “your basket” is never "
        "relabeled “your diet.” Honesty is the feature.",13.5,INK,False,False)]],ls=1.12,anchor=MSO_ANCHOR.MIDDLE)
footer(s,4)
notes(s,"Speaker 1",
      "Think of a new friend who cooks for you. On day one they only know your goal, so they just say what "
      "to aim for. After your first shop they can say what's in your kitchen — but that's shopping, not what "
      "you actually ate, so we call it 'your basket,' not 'your diet.' Once you start tapping meals they get "
      "a feel for it, and after a week they can honestly say 'you're a little low on fibre.' After a month "
      "or two they can show if you're getting better. The rule: if you've only tracked 2 days, we say "
      "'still learning' instead of guessing a wrong number. That honesty is on purpose.")

# ================= SLIDE 5 — KPIs
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Success criteria & KPIs"); heading(s,0.9,0.92,11.5,"How we'll know it works"); rule(s,0.92,1.62)
txt(s,0.9,1.85,11.4,0.4,[[R("Core hypothesis: ",14,INK_SOFT,False,False),
      R("do users trust a receipt-inferred recommendation enough to act — and does the check-off habit stick?",14,INK,True,False)]])
tiles=[("Week 1","Check-off sticks?","the make-or-break — does tapping survive novelty"),
       ("30%","Behaviour shift","adjust their next cart to our suggestion"),
       (">20%","Recommendation acted on","marked “done”, not dismissed unread"),
       ("<2 min","Time to value","from upload to first personalised result")]
tw=2.86;x0=0.9;yy=2.5
for i,(big,lab,sub) in enumerate(tiles):
    stat(s,x0+i*(tw+0.24),yy,tw,2.0,big,lab,sub,accent=SAGE_DK if i==0 else SAGE)
rect(s,0.9,4.95,11.53,1.15,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,1.2,5.12,11,0.85,
    [[R("The check-off habit gates everything.  ",13.5,SAGE_DK,True,False),
      R("Prove the weekly tap sticks first, then layer in the full 10-KPI funnel across behaviour, "
        "business and satisfaction.",13.5,INK,False,False)]],ls=1.12,anchor=MSO_ANCHOR.MIDDLE)
footer(s,5)
notes(s,"Speaker 1",
      "The one number that makes or breaks us: do people keep tapping after week 1? If the daily tap dies, "
      "the whole thing dies — so we test that first. Our north star: at least 3 in 10 people change what "
      "they buy next time, without any manual food diary. Plus: more than 1 in 5 act on our tip, and they "
      "see something useful in under 2 minutes. A bigger 10-KPI scorecard sits behind these.")

# ================= SLIDE 6 — Roadmap (absorbs the old "future" slide)
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Roadmap & milestones"); heading(s,0.9,0.92,11.5,"Where we are — and what's next"); rule(s,0.92,1.62)
ty=2.35; rect(s,1.1,ty+0.98,11.2,0.05,fill=LINE)
cols=[("DONE",SAGE_DK,"Weeks 1–2",["On-device OCR → pantry → gaps → recommendation","Daily check-off log + expiry-aware pantry",
        "Recipes, easy swaps, progress tracking","Nutri-Coach + honest empty states","GDPR export/erasure · DE/EN"]),
      ("IN PROGRESS",SAGE,"Now — Week 3",["1-tap portion — cut check-off friction","Sharper matching (learned Tier-0 + text cleanup)",
        "A/B test: when to ask “what did you eat?”","Live user testing + check-off retention"]),
      ("NEXT",INK_SOFT,"Post-midterm",["Embedding re-ranker to push Recall@5 past 70%","Personalised portion & consumption ML",
        "Instrument the 10-KPI funnel","Weekly / monthly trend views"])]
cw=3.63;x0=1.1
for i,(tag,col,when,items) in enumerate(cols):
    x=x0+i*(cw+0.24)
    rect(s,x+cw/2-0.13,ty+0.86,0.3,0.3,fill=col,shape=MSO_SHAPE.OVAL)
    pill(s,x,ty-0.05,cw*0.55,tag,fill=col)
    txt(s,x+cw*0.58,ty-0.02,cw*0.5,0.44,[[R(when,11.5,INK_SOFT,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    cy=ty+1.5
    rect(s,x,cy,cw,3.15,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    for j,it in enumerate(items):
        yj=cy+0.26+j*0.55
        mark="✓" if i==0 else ("•" if i==2 else "▸")
        txt(s,x+0.3,yj,0.4,0.4,[[R(mark,13,col,True,False)]])
        txt(s,x+0.7,yj-0.02,cw-0.95,0.6,[[R(it,12.5,INK,False,False)]],ls=1.0)
footer(s,6)
notes(s,"Speaker 1",
      "Three weeks of work. DONE: the whole loop already runs end to end. IN PROGRESS: smarter matching and "
      "our first real user tests. NEXT (after today): a smarter matching model to push accuracy past 70%, "
      "personalised portions, and the full success-scorecard. Rather than just tell you it works — let me "
      "show you the live app now.")

# ================= SLIDE 7 — LIVE DEMO
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Live demo — the current build"); heading(s,0.9,0.92,11.5,"See it run, end-to-end, on a real receipt"); rule(s,0.92,1.62)
flow=[("Chat onboarding","goal, needs, exclusions"),
      ("Upload receipt","on-device OCR — private, offline"),
      ("Your Pantry","what's in the house, expiry-aware"),
      ("Daily log","one tap = what you ate"),
      ("Insights","weekly gaps + “use what you have” + Next Cart"),
      ("Nutri-Coach + erase","plain-language why · GDPR erase")]
x0=0.9; yy=2.1
for i,(t,sub) in enumerate(flow):
    y=yy+i*0.72
    rect(s,x0,y,0.5,0.5,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,x0,y-0.02,0.5,0.5,[[R(str(i+1),15,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x0+0.7,y-0.04,5.7,0.4,[[R(t,15,INK,True,False),R("   —  "+sub,11.5,INK_SOFT,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if i<5: rect(s,x0+0.23,y+0.5,0.03,0.22,fill=SAGE_LT)
rx=7.25
rect(s,rx,1.95,5.2,4.55,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
rect(s,rx,1.95,5.2,0.66,fill=SAGE,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.05)
txt(s,rx+0.3,1.99,4.6,0.6,[[R("NutriWise  ·  Your results",13,WHITE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
txt(s,rx+0.3,2.78,4.6,0.5,[[R("Estimates from your groceries + daily log — not medical advice",10.5,INK_SOFT,False,True)]])
rect(s,rx+0.3,3.25,4.6,1.5,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08)
txt(s,rx+0.55,3.4,4.1,0.4,[[R("YOUR NEXT CART",11,SAGE_DK,True,False)]])
txt(s,rx+0.55,3.75,4.1,0.5,[[R("+ Red lentils (dried)",17,INK,True,False)]])
txt(s,rx+0.55,4.2,4.1,0.5,[[R("Closes your fiber & plant-protein gap — and it's peanut-free.",11.5,INK_SOFT,False,False)]],ls=1.05)
for i,(lab,val) in enumerate([("Fiber","low"),("Protein","ok"),("Sugar","ok")]):
    bx=rx+0.3+i*1.55
    rect(s,bx,4.95,1.4,0.85,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.12)
    txt(s,bx,5.05,1.4,0.4,[[R(lab,11,INK_SOFT,True,False)]],align=PP_ALIGN.CENTER)
    txt(s,bx,5.35,1.4,0.4,[[R(val,15,CLAY if val=="low" else SAGE_DK,True,False)]],align=PP_ALIGN.CENTER)
txt(s,rx,6.05,5.2,0.4,[[R("illustrative — shown live in the demo",10,INK_SOFT,False,True)]],align=PP_ALIGN.CENTER)
footer(s,7)
notes(s,"Speaker 1 → hands to Speaker 2",
      "LIVE DEMO, about 2 minutes. Quick tour: chat to set the goal → snap a receipt (say: it reads on the "
      "phone itself, private, works offline, no quota) → land on the Pantry → tap one thing you ate → "
      "Insights shows the week's gaps plus two tips → the coach explains it in plain words, and you can "
      "erase everything anytime. If the coach text ever stops (free daily limit), the numbers and the tip "
      "still work — only the wording gets plainer. Then hand to Speaker 2: 'a tip is only as good as the "
      "data behind it — here's how we made that data trustworthy.'")

# ================= SLIDE 8 — Evaluation metric (gold set + Recall@5)
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Evaluation metric"); heading(s,0.9,0.92,11.7,"The hard part: cryptic receipt text → a real food",size=27); rule(s,0.92,1.62)
txt(s,0.9,1.95,6.2,3.7,
    [[R("“AB Bud 6x0,3l MW”",17,CLAY,True,False),R("   →   ?",17,INK,True,False)],
     [R("",6,INK,False,False)],
     [R("Every receipt line must become a real food with real nutrients — but there is ",14.5,INK,False,False),
      R("no answer key",14.5,INK,True,False),R(" for “what food is this, really?”",14.5,INK,False,False)],
     [R("",5,INK,False,False)],
     [R("So we built one: a ",14.5,INK,False,False),
      R("human-verified gold set",14.5,SAGE_DK,True,False),
      R(" of 180 real receipt lines, each labeled against the ",14.5,INK,False,False),
      R("full 7,140-food corpus",14.5,INK,True,False),R(".",14.5,INK,False,False)],
     [R("Our metric: ",14.5,INK,False,False),
      R("Recall@5",14.5,SAGE_DK,True,False),
      R(" — is the correct food in the top 5 the app shows?",14.5,INK,False,False)]],ls=1.12)
px=7.5; rect(s,px,1.95,4.95,3.85,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05)
txt(s,px+0.35,2.17,4.3,0.4,[[R("WHAT WE MEASURE",12,SAGE_DK,True,False)]])
stat(s,px+0.3,2.62,4.35,1.32,"Recall@5","Right food shortlisted","% where the correct food is in the top 5",accent=SAGE)
stat(s,px+0.3,4.1,4.35,1.32,"3%","Label-error rate","human spot-check of the gold labels",accent=SAGE_DK)
footer(s,8)
notes(s,"Speaker 2",
      "Speaker 2 here. This is not a normal 'train a model' project — the real puzzle is turning messy "
      "receipt text like 'AB Bud 6x0,3l MW' into an actual food. Nobody handed us an answer key, so we made "
      "one: 180 real receipt lines, and for each we found the true match by searching all 7,140 foods — not "
      "just a shortlist, or the score would be fake. A human double-checked the labels; only 3% were off. "
      "Our score: when the app shows its top 5 guesses, is the right food in there? That's Recall@5.")

# ================= SLIDE 9 — Data
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Data — sources, shape & challenges"); heading(s,0.9,0.92,11.5,"Three data worlds, stitched together"); rule(s,0.92,1.62)
srcs=[("Grocery receipts","Real German receipts (Lidl / Netto / Edeka) → on-device OCR (Tesseract) → items. Private: nothing leaves the phone."),
      ("Open Food Facts","Crowd-sourced branded-product data. Fallback for products the BLS structurally lacks."),
      ("BLS 4.0 table","German federal food database — ~7,140 foods. Our primary source of reliable nutrients.")]
cw=3.63;x0=0.9;yy=1.95
for i,(t,b) in enumerate(srcs):
    x=x0+i*(cw+0.24)
    rect(s,x,yy,cw,1.9,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    rect(s,x,yy+0.14,0.08,1.62,fill=SAGE_DK if i==2 else SAGE)
    txt(s,x+0.32,yy+0.2,cw-0.5,0.4,[[R(t,15,INK,True,False)]])
    txt(s,x+0.32,yy+0.72,cw-0.55,1.1,[[R(b,11.5,INK_SOFT,False,False)]],ls=1.1)
rect(s,0.9,4.1,5.55,1.7,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,1.15,4.25,5.1,0.4,[[R("THE GOLD SET (matching test)",11.5,SAGE_DK,True,False)]])
txt(s,1.15,4.67,5.15,1.1,
    [[R("45 → 180",20,SAGE_DK,True,False),R("  real receipt lines   ·   ",12.5,INK,False,False),
      R("152",18,SAGE_DK,True,False),R(" unique",12.5,INK,False,False)],
     [R("Earlier EDA sample: 76 items across 227 purchases.",11,INK_SOFT,False,False)],
     [R("Each line labeled vs. the full 7,140-food corpus.",11,INK_SOFT,False,False)]],ls=1.12,space_after=3)
rect(s,6.7,4.1,5.73,1.7,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05)
txt(s,6.95,4.25,5.3,0.4,[[R("HONEST DATA CHALLENGES",11.5,CLAY,True,False)]])
for i,t in enumerate(["Cryptic OCR, brand-vs-generic (“Budweiser” → “Pilsner Bier”)",
                      "Real receipts are harder than clean text (brand, size, “bio”)",
                      "USDA tested & rejected — blind to German foods (Quark, Schmand)"]):
    txt(s,6.95,4.62+i*0.37,5.35,0.4,[[R("•  ",12,CLAY,True,False),R(t,11.5,INK,False,False)]],ls=1.0)
footer(s,9)
notes(s,"Speaker 2",
      "Three data worlds. One: real German receipts, read on the phone itself — private, no cloud. Two: "
      "Open Food Facts, a big crowd list of branded products, as a backup. Three: the official German food "
      "table, 7,140 foods, our trusted source for nutrients. For testing we grew our answer key from 45 to "
      "180 real receipt lines. Honest hurdles: receipt text is cryptic; real receipts are harder than tidy "
      "text; and we tried the US food database but dropped it — it simply doesn't know German foods like "
      "Quark or Schmand.")

# ================= SLIDE 10 — EDA
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Exploratory data analysis"); heading(s,0.9,0.92,11.5,"What real baskets actually look like"); rule(s,0.92,1.62)
rect(s,0.9,1.95,5.75,4.3,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
s.shapes.add_picture(CH("eda_categories.png"),Inches(1.05),Inches(2.12),height=Inches(3.9))
rect(s,6.85,1.95,5.58,4.3,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
s.shapes.add_picture(CH("eda_tail.png"),Inches(7.0),Inches(2.12),height=Inches(3.9))
txt(s,0.9,6.32,11.5,0.5,
    [[R("Takeaway:  ",12.5,SAGE_DK,True,False),
      R("baskets are produce- and dairy-heavy, but the long tail is brutal — half the items appear only "
        "once, so the matcher must generalise, not memorise.",12.5,INK,False,False)]],ls=1.05)
footer(s,10)
notes(s,"Speaker 2",
      "Two quick pictures. Left: most shopping is fruit, veg and dairy — the basics. Right: but there's a "
      "long tail — half the items show up only once. So we can't just memorise the popular foods; the "
      "matcher has to handle weird, one-off, badly-abbreviated names it has never seen. That shaped our "
      "whole design.")

# ================= SLIDE 11 — Pipeline (matching resolver)
s=slide(); contentbg(s); kicker(s,0.9,0.6,"ML workflow / data pipeline"); heading(s,0.9,0.92,11.7,"A tiered matcher — cheap & certain first, fallback last"); rule(s,0.92,1.62)
stages=[("Receipt","image / text"),("Extract","on-device OCR → items"),("Clean","strip qty · price · brand")]
tiers=[("Tier 0","Learned verified matches",SAGE_DK,"human-confirmed · always trusted"),
       ("Tier 1","BLS 4.0 — German nutrients",SAGE,"primary source · 174 of 180 lines"),
       ("Tier 2","Open Food Facts fallback",SAGE_LT,"branded gaps · 4 of 180 lines"),
       ("open","Left honest / unresolved",SAND,"2 ambiguous lines — not guessed")]
yy=2.05;bw=2.55;x0=0.9
for i,(t,sub) in enumerate(stages):
    x=x0+i*(bw+0.55)
    rect(s,x,yy,bw,1.0,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08)
    txt(s,x,yy+0.18,bw,0.4,[[R(t,15,WHITE,True,False)]],align=PP_ALIGN.CENTER)
    txt(s,x,yy+0.6,bw,0.35,[[R(sub,11,SAGE_LT,False,False)]],align=PP_ALIGN.CENTER)
    if i<2: txt(s,x+bw+0.04,yy+0.2,0.5,0.6,[[R("→",24,SAGE,True,False)]])
txt(s,x0+3*(bw+0.55)-0.35,yy+0.2,1.4,0.6,[[R("→",24,SAGE,True,False)]])
ty=3.35
for i,(tag,name,col,sub) in enumerate(tiers):
    y=ty+i*0.72
    rect(s,0.9,y,8.4,0.62,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.12)
    rect(s,0.9,y,1.55,0.62,fill=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.12)
    tc=INK if col in (SAGE_LT,SAND) else WHITE
    txt(s,0.9,y,1.55,0.62,[[R(tag,13,tc,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,2.65,y,3.6,0.62,[[R(name,13.5,INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,6.0,y,3.2,0.62,[[R(sub,11.5,INK_SOFT,False,True)]],anchor=MSO_ANCHOR.MIDDLE)
rect(s,9.65,3.35,2.78,2.78,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,9.85,3.55,2.4,0.4,[[R("OUTPUT",11.5,SAGE_DK,True,False)]])
for i,t in enumerate(["Nutrients (macros + micros)","Confidence label","Data-source citation","→ Gaps + Next Cart"]):
    txt(s,9.85,4.0+i*0.53,2.5,0.5,[[R("• ",12,SAGE,True,False),R(t,12,INK,False,False)]],ls=1.0)
txt(s,0.9,6.35,8.4,0.4,[[R("Rule-based on purpose — the embedding only breaks ties, so the logic stays transparent for a health MVP.",12,INK_SOFT,False,True)]])
footer(s,11)
notes(s,"Speaker 2",
      "One important note first: these matching tiers are about linking receipt text to a food — different "
      "from the honesty tiers earlier, which were about the score. Here it's a ladder: read the receipt, "
      "clean it up, then try the safest match first. Tier 0: matches a human already confirmed. Tier 1: the "
      "German food table — our main source, it handled 174 of 180 lines. Tier 2: Open Food Facts for the "
      "few branded gaps. And if it's genuinely unclear (2 lines), we leave it open instead of guessing. "
      "It's rules on purpose — trustworthy for a health app — and the smart 'embedding' only breaks ties.")

# ================= SLIDE 12 — Results (baseline & performance, ML Step 2)
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Baseline model & performance"); heading(s,0.9,0.92,11.7,"Measured honestly on real receipts"); rule(s,0.92,1.62)
rect(s,0.9,1.9,7.05,3.75,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
s.shapes.add_picture(CH("cleanup.png"),Inches(1.15),Inches(2.1),height=Inches(3.35))
stat(s,8.25,1.9,4.18,1.15,"99%","Lines resolved","174 BLS · 4 OFF · 2 left open, of 180",accent=SAGE)
stat(s,8.25,3.15,4.18,1.15,"47%","Recall@5 (cleaned)","right food in top 5 — up from 32% raw",accent=SAGE_DK)
stat(s,8.25,4.4,4.18,1.15,"56% → 47%","Real is harder","clean 45 vs. real receipts (Recall@5)",accent=SAGE)
rect(s,0.9,5.82,11.53,0.85,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,1.2,5.93,11,0.7,
    [[R("The honest finding:  ",13,SAGE_DK,True,False),
      R("on raw text, simple fuzzy beats the “smart” hybrid — so the embedding is a tie-breaker, not the "
        "main matcher. Cleaning the text is the real lever; the learned Tier-0 store closes the rest.",
        13,INK,False,False)]],ls=1.08,anchor=MSO_ANCHOR.MIDDLE)
footer(s,12)
notes(s,"Speaker 2",
      "Our honest scorecard. First, we found a match for 99% of lines — 174 from the German table, 4 from "
      "Open Food Facts, and 2 we left open on purpose. Second, the big lever is cleaning the receipt text: "
      "the right food lands in the top 5 about 47% of the time after cleanup, up from 32% on raw text. "
      "Third, honesty: real receipts are harder than the tidy 45 we started with (56% → 47%), and a plain "
      "'fuzzy' search actually beats our fancy one on messy text — so the AI part only breaks ties for now. "
      "The plan to close the gap: every user correction is saved and reused. (Backup slide has our earlier "
      "pilot numbers if asked.)")

# ================= SLIDE 13 — Thank you
s=slide(); bg(s,CANVAS)
rect(s,0,0,13.333,0.14,fill=SAGE)
rect(s,0,7.36,13.333,0.14,fill=SAGE)
logo(s,0.9,0.9)
txt(s,0.9,2.5,11.5,1.2,[[R("Thank you.",50,INK,True,False)]])
rule(s,0.92,3.7,3.4)
txt(s,0.92,3.95,11,0.6,[[R("Upload a grocery receipt. Get one smarter thing to buy next time.",18,SAGE_DK,False,True)]])
txt(s,0.9,4.95,11.5,1.0,
    [[R("Jennifer Rake",15,INK,True,False),R("   &   ",15,INK_SOFT,False,False),R("Stuart Kasemeier",15,INK,True,False)],
     [R("NutriWise  ·  AI PM Capstone  ·  Midterm V2  —  Questions & feedback welcome",12.5,INK_SOFT,False,False)]],space_after=8)
notes(s,"Both",
      "Close on the one-liner and invite questions. If pushed on accuracy: we measure honestly on real "
      "receipts (Recall@5 47% after cleanup, 99% of lines resolved), the clear next lever is the learned "
      "match store, and the deeper question — do people trust it and keep tapping — is what live user "
      "testing answers next.")

# ================= SLIDE 14 — BACKUP: earlier pilot (OFF vs BLS, LLM-judged)
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Backup — earlier pilot measurement")
heading(s,0.9,0.92,11.7,"Pilot: OFF vs BLS, judged head-to-head",size=26); rule(s,0.92,1.62)
txt(s,0.9,1.68,11.5,0.4,[[R("Different method than Recall@5 — an earlier LLM-judge pilot on 76 items. Not directly comparable; kept for context.",12,INK_SOFT,False,True)]])
rect(s,0.9,2.15,7.55,3.5,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
s.shapes.add_picture(CH("baseline.png"),Inches(1.05),Inches(2.35),width=Inches(7.25))
stat(s,8.75,2.15,3.68,1.05,"66%","OFF usable (pilot)","correct + partial, LLM-judged",accent=SAGE)
stat(s,8.75,3.32,3.68,1.05,"2×","OFF beats BLS","54% vs 32% correct on identity",accent=SAGE_DK)
stat(s,8.75,4.49,3.68,1.05,"96%","Judge agreement","3/3 unanimous — pilot metric",accent=SAGE)
txt(s,0.9,5.85,11.5,0.5,[[R("Why it still matters:  ",12.5,SAGE_DK,True,False),
    R("it's why BLS leads and OFF is only a fallback — and why we moved to the stricter, human-verified Recall@5.",12.5,INK,False,False)]],ls=1.05)
notes(s,"Speaker 2 — only if asked",
      "Backup for questions. This was our earlier pilot on 76 items, judged by an LLM panel — a different "
      "method, so don't compare its 66% to the 47% Recall@5. What it told us: Open Food Facts wins on "
      "identity, the German table matched worse on its own — which is exactly why the table leads for "
      "nutrients only, and why we then built the stricter, human-checked Recall@5 test.")

out=os.path.join(OUTDIR,"NutriWise_Midterm_Presentation_v2.pptx")
prs.save(out); print("saved",out,"slides:",len(prs.slides._sldIdLst))
