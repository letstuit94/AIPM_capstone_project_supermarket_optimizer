#!/usr/bin/env python3
"""NutriWise Midterm-Deck V1 — DEUTSCH. Corporate Design (Salbeigrün, sans-serif).
Charts + .pptx, 14 Folien inkl. Live-Demo. Alle Zahlen aus dem Repo."""
import json, collections, os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------- corporate palette
SAGE     = "7C9A6A"; SAGE_DK  = "5F7B4F"; SAGE_LT  = "B7C9A9"; SAGE_SOFT= "EEF2EA"
CLAY     = "B36A4A"; SAND     = "CFC9BA"; INK      = "1D1D21"; INK_SOFT = "6E6F74"
CANVAS   = "FAFAFA"; SURFACE  = "FFFFFF"; LINE     = "E5E4E0"; WHITE    = "FFFFFF"
FONT     = "Arial"

def C(h): return RGBColor.from_string(h)
SCRATCH = os.path.dirname(os.path.abspath(__file__))
ROOT    = "/Users/jenniferrake/Desktop/GitHub_Bootcamp/AIPM_capstone_project_supermarket_optimizer"
OUTDIR  = os.path.join(ROOT, "midterm_presentation")
CH = lambda n: os.path.join(SCRATCH, n)

# ---------------------------------------------------- data
d = json.load(open(os.path.join(ROOT, "bls_off_judgments.json")))
N = len(d)
off = collections.Counter(x["off_verdict"] for x in d)
bls = collections.Counter(x["bls_verdict"] for x in d)
GROUP = {
    "Gemüse":"Obst & Gemüse","Obst":"Obst & Gemüse","Obst & Gemüse":"Obst & Gemüse",
    "Milchprodukte":"Milch, Käse & Eier","Käse":"Milch, Käse & Eier",
    "Milchprodukte & Alternativen":"Milch, Käse & Eier","Eier":"Milch, Käse & Eier",
    "Backwaren":"Brot & Getreide","Müsli & Cerealien":"Brot & Getreide",
    "Nudeln & Reis":"Brot & Getreide","grain":"Brot & Getreide",
    "Fleisch":"Fleisch & Fisch","Fleisch/Wurst":"Fleisch & Fisch","Fisch":"Fleisch & Fisch",
    "Hülsenfrüchte":"Hülsenfr. & Protein","protein":"Hülsenfr. & Protein",
    "Getränke":"Getränke & Tee","Tee":"Getränke & Tee",
}
grp = collections.Counter()
for x in d:
    grp[GROUP.get(x["category"], "Sonstiges / verarbeitet")] += 1

# ---------------------------------------------------- matplotlib defaults
plt.rcParams.update({
    "font.family":"sans-serif","font.sans-serif":["Arial","Helvetica","DejaVu Sans"],
    "font.size":15,"text.color":"#"+INK,"axes.edgecolor":"#"+LINE,
    "axes.labelcolor":"#"+INK,"xtick.color":"#"+INK_SOFT,"ytick.color":"#"+INK,
    "figure.facecolor":"#"+SURFACE,"axes.facecolor":"#"+SURFACE,"savefig.facecolor":"#"+SURFACE,
})
def strip(ax, keep_x=False):
    for s in ["top","right","left"]: ax.spines[s].set_visible(False)
    if not keep_x:
        ax.spines["bottom"].set_visible(False); ax.tick_params(bottom=False)
    ax.tick_params(left=False)

# ---- CHART 1
labels=[k for k,_ in grp.most_common()]; vals=[v for _,v in grp.most_common()]
fig,ax=plt.subplots(figsize=(6.6,4.3),dpi=200)
y=list(range(len(labels)))[::-1]
bars=ax.barh(y,vals,color="#"+SAGE,height=0.66); bars[0].set_color("#"+SAGE_DK)
ax.set_yticks(y); ax.set_yticklabels(labels,fontsize=13); ax.set_xticks([])
for yi,v in zip(y,vals):
    ax.text(v+0.3,yi,str(v),va="center",ha="left",fontsize=13,color="#"+INK,fontweight="bold")
strip(ax); ax.set_xlim(0,max(vals)+3)
ax.set_title("Was Menschen wirklich kaufen — 76 einzigartige Artikel",
             fontsize=14.5,fontweight="bold",color="#"+INK,loc="left",pad=12)
plt.tight_layout(); plt.savefig(CH("eda_categories_de.png"),bbox_inches="tight"); plt.close()

# ---- CHART 2
occ=collections.Counter(x["occurrences"] for x in d)
buckets={"1×":occ[1],"2×":occ[2],"3–4×":occ[3]+occ[4],
         "5–9×":sum(occ[k] for k in range(5,10)),
         "10×+":sum(v for k,v in occ.items() if k>=10)}
fig,ax=plt.subplots(figsize=(6.6,4.3),dpi=200)
bl=list(buckets); bv=list(buckets.values())
bars=ax.bar(bl,bv,color="#"+SAGE_LT,width=0.62); bars[0].set_color("#"+SAGE)
for b,v in zip(bars,bv):
    ax.text(b.get_x()+b.get_width()/2,v+0.4,str(v),ha="center",va="bottom",
            fontsize=13,fontweight="bold",color="#"+INK)
strip(ax,keep_x=True); ax.set_yticks([]); ax.set_ylim(0,max(bv)+4)
ax.tick_params(axis="x",labelsize=13)
ax.set_title("Echte Warenkörbe sind long-tail: die Hälfte aller\nArtikel erscheint nur einmal",
             fontsize=14.5,fontweight="bold",color="#"+INK,loc="left",pad=12)
plt.tight_layout(); plt.savefig(CH("eda_tail_de.png"),bbox_inches="tight"); plt.close()

# ---- CHART 3
order=["correct","partially_correct","no_match","incorrect"]
oleg={"correct":"korrekt","partially_correct":"teilweise","no_match":"kein Treffer","incorrect":"falsch"}
colors={"correct":SAGE,"partially_correct":SAGE_LT,"no_match":SAND,"incorrect":CLAY}
rows=[("OpenFoodFacts\nIdentitäts-Match",off),("BLS-Tabelle\nMatch",bls)]
fig,ax=plt.subplots(figsize=(9.2,3.5),dpi=200); ypos=[1,0]
for (name,cnt),yp in zip(rows,ypos):
    left=0
    for k in order:
        w=cnt.get(k,0)/N*100
        ax.barh(yp,w,left=left,color="#"+colors[k],height=0.5,edgecolor="#"+SURFACE,linewidth=1.5)
        if w>=7:
            ax.text(left+w/2,yp,f"{round(w)}%",ha="center",va="center",fontsize=13,
                    fontweight="bold",color="#"+WHITE if k in("correct","incorrect") else "#"+INK)
        left+=w
ax.set_yticks(ypos); ax.set_yticklabels([r[0] for r in rows],fontsize=13,fontweight="bold")
ax.set_xticks([]); ax.set_xlim(0,100); ax.set_ylim(-0.6,1.6); strip(ax)
handles=[plt.Rectangle((0,0),1,1,color="#"+colors[k]) for k in order]
ax.legend(handles,[oleg[k] for k in order],ncol=4,loc="upper center",
          bbox_to_anchor=(0.5,-0.08),frameon=False,fontsize=12,handlelength=1.1)
plt.tight_layout(); plt.savefig(CH("baseline_de.png"),bbox_inches="tight"); plt.close()
print("charts done")

# ==================================================== PPTX
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def bg(s,h): s.background.fill.solid(); s.background.fill.fore_color.rgb=C(h)
def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,shadow=False,adj=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=C(line); sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if adj is not None:
        try: sp.adjustments[0]=adj
        except Exception: pass
    if shadow:
        from pptx.oxml.ns import qn
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{})
        sh=el.makeelement(qn('a:outerShdw'),{'blurRad':'80000','dist':'30000','dir':'5400000','rotWithShape':'0'})
        clr=el.makeelement(qn('a:srgbClr'),{'val':'1D1D21'}); alp=el.makeelement(qn('a:alpha'),{'val':'12000'})
        clr.append(alp); sh.append(clr); ef.append(sh); el.append(ef)
    return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space_after=4,ls=1.0,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,col,b,it) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it
            r.font.color.rgb=C(col); r.font.name=FONT
    return tb
def R(t,sz,col=INK,b=False,it=False): return (t,sz,col,b,it)
def notes(s,speaker,body): s.notes_slide.notes_text_frame.text=f"[{speaker}]\n{body}"
def kicker(s,x,y,t): txt(s,x,y,8,0.3,[[R(t.upper(),12.5,SAGE_DK,True,False)]])
def heading(s,x,y,w,t,size=29): txt(s,x,y,w,1.0,[[R(t,size,INK,True,False)]],ls=1.0)
def rule(s,x,y,w=1.5): rect(s,x,y,w,0.045,fill=SAGE)
def footer(s,page):
    txt(s,0.6,7.03,5,0.3,[[R("NutriWise",9,INK_SOFT,True,False),R("  ·  V1",9,INK_SOFT,False,False)]])
    txt(s,10.8,7.03,1.9,0.3,[[R(f"{page:02d} / 14",9,INK_SOFT,False,False)]],align=PP_ALIGN.RIGHT)
def logo(s,x,y,dark=False):
    rect(s,x,y,0.34,0.34,fill=SAGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.3)
    txt(s,x+0.46,y-0.03,4,0.42,[[R("NUTRIWISE",13,WHITE if dark else INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
def contentbg(s): bg(s,CANVAS)
def stat(s,x,y,w,h,big,label,sub=None,accent=SAGE):
    rect(s,x,y,w,h,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.06)
    rect(s,x,y+0.14,0.08,h-0.28,fill=accent)
    body=[[R(big,29,accent,True,False)],[R(label,12.5,INK,True,False)]]
    if sub: body.append([R(sub,10.5,INK_SOFT,False,False)])
    txt(s,x+0.3,y+0.14,w-0.42,h-0.26,body,anchor=MSO_ANCHOR.MIDDLE,space_after=3,ls=1.02)
def pill(s,x,y,w,t,fill=SAGE,col=WHITE):
    rect(s,x,y,w,0.44,fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    txt(s,x,y+0.02,w,0.4,[[R(t,11,col,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ================= FOLIE 1 — Titel
s=slide(); bg(s,CANVAS)
rect(s,0,0,13.333,0.14,fill=SAGE)
rect(s,8.9,0.14,4.433,7.36,fill=SAGE_SOFT)
logo(s,0.9,0.85)
txt(s,0.9,2.15,8.0,2.6,
    [[R("Lade deinen Kassenbon hoch.",38,INK,True,False)],
     [R("Bekomm eine klügere Sache",38,SAGE_DK,True,False)],
     [R("für den nächsten Einkauf.",38,SAGE_DK,True,False)]],ls=1.05)
rule(s,0.92,4.55,3.4)
txt(s,0.92,4.8,7.6,0.7,[[R("Ein Kassenbon-basierter Ernährungsassistent — dein Vorrat arbeitet für dich: abhaken statt tippen.",15,INK_SOFT,False,True)]])
txt(s,0.9,5.75,7.8,1.2,
    [[R("Jennifer Rake",15,INK,True,False),R("  · Sprecher 2      ",13,SAGE_DK,True,False),
      R("Stuart Kasemeier",15,INK,True,False),R("  · Sprecher 1",13,INK_SOFT,True,False)],
     [R("AI Product Management Capstone  ·  Midterm-Präsentation  ·  Version 1",12.5,INK_SOFT,False,False)]],space_after=8)
for i,t in enumerate(["Onboarding","Vorrat","Abhaken","Insights"]):
    yy=2.35+i*0.92
    rect(s,9.5,yy,0.26,0.26,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,9.95,yy-0.06,2.9,0.4,[[R(t,14,INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if i<3: rect(s,9.62,yy+0.3,0.02,0.6,fill=SAGE_LT)
notes(s,"Beide — Sprecher 1 beginnt",
      "Sprecher 1 (Stuart) eröffnet mit dem Einzeiler und stellt beide vor. Erwartung setzen: ~10–12 Min., "
      "high-level, plus kurze Live-Demo. Das ist V1 unserer finalen Präsentation. Sprecher 2 (Jennifer) "
      "übernimmt nach der Demo den Daten-/Modell-Teil.")

# ================= FOLIE 2 — Problem
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Das Problem"); heading(s,0.9,0.92,11,"Gesund essen sollte keine Tabelle erfordern"); rule(s,0.92,1.62)
txt(s,0.9,2.0,6.1,3.6,
    [[R("Gesundheitsbewusste Erwachsene (25–45) wollen besser essen — aber jedes Tool verlangt, ",15.5,INK,False,False),
      R("jede Mahlzeit von Hand zu erfassen.",15.5,INK,True,False)],
     [R("",7,INK,False,False)],
     [R("Diese Gewohnheit stirbt in Wochen. Man gibt auf — oder rät.",15.5,INK,False,False)],
     [R("",7,INK,False,False)],
     [R("Dabei existiert der eine ehrliche Beleg dessen, was wirklich gegessen wird, längst — ",15.5,INK,False,False),
      R("der Kassenbon.",15.5,SAGE_DK,True,False)]],ls=1.14)
px=7.4; rect(s,px,1.95,5.05,3.85,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05)
txt(s,px+0.35,2.2,4.4,0.4,[[R("WARUM ES ZÄHLT",12,SAGE_DK,True,False)]])
for i,t in enumerate(["Manuelles Erfassen ist Grund #1, warum Ernährungs-Apps abgebrochen werden",
                      "Kassenbons sind passiv — null Zusatzaufwand für den User",
                      "Eine klare „nächster Einkauf“-Entscheidung schlägt ein weiteres Dashboard"]):
    yy=2.68+i*1.0
    rect(s,px+0.35,yy+0.03,0.34,0.34,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,px+0.35,yy-0.01,0.34,0.34,[[R(str(i+1),13,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,px+0.85,yy-0.05,3.9,0.9,[[R(t,13.5,INK,False,False)]],ls=1.05)
footer(s,2)
notes(s,"Sprecher 1",
      "Den Schmerz rahmen: Menschen wollen besser essen, aber manuelles Loggen (à la MyFitnessPal) ist "
      "anstrengend und wird abgebrochen. Zielgruppe = gesundheitsbewusst, 25–45. Kern-Insight: der Bon ist "
      "ein passiver, ehrlicher Proxy für das tatsächlich Gekaufte — kein Loggen nötig. Das ist unser Wedge.")

# ================= FOLIE 3 — Lösung
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Die Lösung"); heading(s,0.9,0.92,11.5,"Ein Bon rein → ein Vorrat, der für dich arbeitet"); rule(s,0.92,1.62)
steps=[("1","Upload","Bon fotografieren/einfügen — On-Device-OCR baut deinen Vorrat. Privat, keine API, offline."),
       ("2","Dein Vorrat","Was im Haus ist, nach Frische sortiert — kein manuelles Erfassen."),
       ("3","Abhaken","Ein Tap erfasst, was du wirklich gegessen hast — dein echtes Ist-Signal."),
       ("4","Lücken + Tipp","Wochen-Lücken → „nutze, was du hast“ + ein Next-Cart-Kauf, je mit klarer Begründung.")]
cw=2.86;gp=0.24;x0=0.9;yy=2.1
for i,(n,t,body) in enumerate(steps):
    x=x0+i*(cw+gp)
    rect(s,x,yy,cw,3.25,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    rect(s,x,yy,cw,0.72,fill=SAGE if i<3 else SAGE_DK,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.08)
    txt(s,x,yy+0.08,cw,0.6,[[R("SCHRITT "+n,12,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.25,yy+0.92,cw-0.5,0.5,[[R(t,17,INK,True,False)]])
    txt(s,x+0.25,yy+1.52,cw-0.5,1.6,[[R(body,12.5,INK_SOFT,False,False)]],ls=1.08)
    if i<3: txt(s,x+cw-0.02,yy+1.2,0.34,0.5,[[R("›",26,SAGE,True,False)]])
txt(s,0.9,5.65,11.5,0.7,
    [[R("Alles ist gekennzeichnet als ",13.5,INK,False,True),
      R("„Schätzung aus Einkauf + deinen Taps.“",13.5,SAGE_DK,True,True),
      R("  Keine medizinische Beratung.",13.5,INK,False,True)]])
footer(s,3)
notes(s,"Sprecher 1",
      "Den echten Loop zeigen: Upload baut den Vorrat → du tappst, was du gegessen hast → Wochen-Lücken → "
      "zwei Empfehlungen („nutze, was du hast“ + ein Next-Cart-Kauf). Die Produktwette: ein lebendiger "
      "Vorrat plus eine klare Entscheidung schlägt ein Ernährungstagebuch. Trust-Framing: Schätzung aus "
      "Einkauf + deinen Taps, keine medizinische Beratung.")

# ================= FOLIE 4 — UI / User-Journey
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Das Produkt heute"); heading(s,0.9,0.92,11.7,"Vom Bon zur täglichen Gewohnheit — die User-Journey",size=26); rule(s,0.92,1.62)
pw=2.55; ph=3.5; py=2.1; xs=[0.95,3.92,6.89,9.86]
caps=[("1 · Onboarding-Chat","dialogbasiertes Profil"),("2 · Dein Vorrat","Home · Ablauf-bewusst"),
      ("3 · Tages-Log","1 Tap = 1 Portion"),("4 · Insights","Wochen-Lücken + Tipps")]
def mpshell(x,title):
    rect(s,x,py,pw,ph,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.07)
    rect(s,x,py,pw,0.5,fill=SAGE,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.12)
    txt(s,x+0.2,py+0.05,pw-0.3,0.42,[[R(title,10,WHITE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
for k in range(3):
    txt(s,xs[k]+pw+0.02,py+1.45,0.4,0.5,[[R("›",22,SAGE,True,False)]],align=PP_ALIGN.CENTER)
for k,(t,sub) in enumerate(caps):
    txt(s,xs[k],py+ph+0.12,pw,0.55,[[R(t,11.5,INK,True,False)],[R(sub,10,INK_SOFT,False,False)]],align=PP_ALIGN.CENTER,space_after=1)
# phone 1
x=xs[0]; mpshell(x,"NutriWise")
rect(s,x+0.2,py+0.72,1.55,0.4,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.45)
txt(s,x+0.33,py+0.74,1.4,0.36,[[R("Dein Ziel?",8,INK,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
rect(s,x+0.8,py+1.24,1.55,0.4,fill=SAGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.45)
txt(s,x+0.93,py+1.26,1.4,0.36,[[R("Gesünder essen",7.5,WHITE,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
for j,lab in enumerate(["Wenig","Mittel","Hoch"]):
    cxp=x+0.2+j*0.72
    rect(s,cxp,py+1.92,0.66,0.34,fill=SURFACE,line=SAGE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    txt(s,cxp,py+1.94,0.66,0.3,[[R(lab,7,SAGE_DK,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
for j in range(4):
    rect(s,x+0.2+j*0.2,py+2.6,0.14,0.05,fill=SAGE if j==0 else LINE)
# phone 2
x=xs[1]; mpshell(x,"Dein Vorrat")
rect(s,x+0.2,py+0.68,2.15,0.55,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.2)
txt(s,x+0.33,py+0.72,1.95,0.5,[[R("EMPFOHLEN",7,SAGE_DK,True,False)],[R("🥛 Trink deine Buttermilch",8.5,INK,True,False)]],space_after=1)
for j,(nm,bd,cc) in enumerate([("🍗 Hähnchenbrust","1 Tag",CLAY),("🥛 Buttermilch","2 Tage",CLAY),("🥦 Brokkoli","frisch",SAGE_DK)]):
    yb=py+1.42+j*0.6
    rect(s,x+0.2,yb,2.15,0.5,fill=SURFACE,line=LINE,lw=0.6,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.2)
    txt(s,x+0.33,yb,1.5,0.5,[[R(nm,8,INK,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+1.7,yb,0.6,0.5,[[R(bd,7,cc,True,False)]],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.RIGHT)
# phone 3
x=xs[2]; mpshell(x,"Tages-Log")
for j,(nm,val,frac,cc) in enumerate([("Protein","62 / 120 g",0.52,CLAY),("Ballaststoffe","11 / 30 g",0.37,CLAY),("Kalorien","1.340 / 2.100",0.64,SAGE)]):
    yb=py+0.82+j*0.78
    txt(s,x+0.2,yb,2.15,0.26,[[R(nm,7.5,INK,True,False),R("   "+val,7,INK_SOFT,False,False)]])
    rect(s,x+0.2,yb+0.3,2.15,0.14,fill=LINE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    rect(s,x+0.2,yb+0.3,2.15*frac,0.14,fill=cc,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
# phone 4
x=xs[3]; mpshell(x,"Insights")
txt(s,x+0.2,py+0.66,2.15,0.66,[[R("72",30,SAGE_DK,True,False)]],align=PP_ALIGN.CENTER)
txt(s,x+0.2,py+1.32,2.15,0.3,[[R("Wochen-Score · Gut",8,INK_SOFT,False,False)]],align=PP_ALIGN.CENTER)
for j,(nm,st,cc) in enumerate([("Protein","niedrig",CLAY),("Ballastst.","niedrig",CLAY),("Zucker","ok",SAGE_DK),("Kalorien","ok",SAGE_DK)]):
    cxp=x+0.2+(j%2)*1.1; cyp=py+1.72+(j//2)*0.62
    rect(s,cxp,cyp,1.02,0.5,fill=SURFACE,line=LINE,lw=0.6,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.2)
    txt(s,cxp+0.11,cyp,0.85,0.5,[[R(nm,7,INK,False,False)],[R(st,7.5,cc,True,False)]],space_after=0,anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.9,6.88,11.5,0.3,[[R("Voll klickbarer Prototyp — Salbeigrün-Corporate-Design, hell & dunkel.",11,INK_SOFT,False,True)]],align=PP_ALIGN.CENTER)
footer(s,4)
notes(s,"Sprecher 1 oder 2",
      "Die echte Produktform auf einen Blick: dialogbasiertes Onboarding setzt die Zielwerte; der Vorrat "
      "ist der Startbildschirm, zu dem man zurückkehrt (nach Ablauf sortiert); ein Tap im Tages-Log erfasst "
      "Gegessenes; Insights macht aus einer Woche Taps Lücken und zwei Empfehlungen. Ein voll klickbarer "
      "Prototyp existiert, falls das Publikum es live sehen will.")

# ================= FOLIE 5 — KPIs
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Erfolgskriterien & KPIs"); heading(s,0.9,0.92,11.5,"Woran wir Erfolg messen"); rule(s,0.92,1.62)
txt(s,0.9,1.85,11.5,0.4,[[R("Kernhypothese: ",14,INK_SOFT,False,False),
      R("Vertrauen Nutzer einer aus dem Bon abgeleiteten Empfehlung genug, um zu handeln — und klebt das Abhaken?",14,INK,True,False)]])
tiles=[("Woche 1","Abhaken klebt?","Make-or-Break — überlebt das Tippen die Neuheit"),
       ("30%","Verhalten ändert sich","passen den nächsten Einkauf an unseren Tipp an"),
       (">20%","Empfehlung befolgt","als „erledigt“ markiert, nicht ungelesen verworfen"),
       ("<2 Min","Zeit bis Nutzen","vom Upload zum ersten personalisierten Ergebnis")]
tw=2.86;x0=0.9;yy=2.45
for i,(big,lab,sub) in enumerate(tiles):
    stat(s,x0+i*(tw+0.24),yy,tw,2.0,big,lab,sub,accent=SAGE_DK if i==0 else SAGE)
rect(s,0.9,4.85,11.53,1.35,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,1.2,5.02,11,1.05,
    [[R("Das Abhaken entscheidet über alles.  ",13.5,SAGE_DK,True,False),
      R("Überlebt das Tippen Woche 1 nicht, veralten Vorrat und Empfehlungen — also beweisen wir zuerst "
        "diesen Loop und legen dann den vollen 10-KPI-Funnel über Verhalten, Business und Zufriedenheit.",13.5,INK,False,False)]],
    ls=1.15,anchor=MSO_ANCHOR.MIDDLE)
footer(s,5)
notes(s,"Sprecher 1",
      "Der Make-or-Break-Frühindikator ist die Abhak-Retention — überlebt der tägliche Tap die Neuheit? "
      "Wenn nicht, stirbt das gesamte Vorratsmodell, also beweisen wir das zuerst. North-Star: ≥30% ändern "
      "ihren nächsten Einkauf ohne manuelles Loggen. Ergänzend: Aktionsrate >20%, Zeit bis Nutzen <2 Min. "
      "Dahinter das volle 10-KPI-Framework.")

# ================= FOLIE 6 — Roadmap
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Roadmap & Meilensteine"); heading(s,0.9,0.92,11.5,"Wo wir stehen"); rule(s,0.92,1.62)
ty=2.35; rect(s,1.1,ty+0.98,11.2,0.05,fill=LINE)
cols=[("FERTIG",SAGE_DK,"Woche 1–2",["On-Device-OCR → Vorrat → Lücken → Empfehlung","Tages-Log (Abhaken) + Ablauf-bewusster Vorrat",
        "Rezepte, Easy Swaps, Progress-Tracking","Nutri-Coach + transparente Leerzustände","DSGVO-Export/Löschung · Engagement-Banner · DE/EN"]),
      ("IN ARBEIT",SAGE,"Jetzt — Woche 3",["1-Tap-Portion — Reibung senken","Besseres Matching (gelernter Tier-0 + Normalisierung)",
        "A/B-Test: wann nach „was gegessen?“ fragen","Live-User-Tests + Abhak-Retention"]),
      ("NÄCHSTES",INK_SOFT,"Nach Midterm",["Advanced-Matching (Embeddings) über 70 %","10-KPI-Funnel instrumentieren",
        "Wochen-/Monats-Trends"])]
cw=3.63;x0=1.1
for i,(tag,col,when,items) in enumerate(cols):
    x=x0+i*(cw+0.24)
    rect(s,x+cw/2-0.13,ty+0.86,0.3,0.3,fill=col,shape=MSO_SHAPE.OVAL)
    pill(s,x,ty-0.05,cw*0.55,tag,fill=col)
    txt(s,x+cw*0.58,ty-0.02,cw*0.5,0.44,[[R(when,11.5,INK_SOFT,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    cy=ty+1.5
    rect(s,x,cy,cw,3.05,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    for j,it in enumerate(items):
        yy=cy+0.26+j*0.53
        mark="✓" if i==0 else ("•" if i==2 else "▸")
        txt(s,x+0.3,yy,0.4,0.4,[[R(mark,13,col,True,False)]])
        txt(s,x+0.7,yy-0.02,cw-0.95,0.6,[[R(it,12.5,INK,False,False)]],ls=1.0)
footer(s,6)
notes(s,"Sprecher 1",
      "Drei-Wochen-Sprint. FERTIG: der ganze Kern-Loop läuft end-to-end. IN ARBEIT: besseres Matching + "
      "erste 5 User-Tests + 1-Tap-Portion. NÄCHSTES: Genauigkeit über 70%, KPI-Instrumentierung. Und statt "
      "es nur zu erzählen — ich zeige euch den aktuellen Stand.")

# ================= FOLIE 7 — Live-Demo
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Live-Demo — der aktuelle Stand"); heading(s,0.9,0.92,11.5,"Einmal komplett durch — an einem echten Bon"); rule(s,0.92,1.62)
flow=[("Onboarding-Chat","dialogbasiertes Profil — Ziel, Bedarf, Ausschlüsse"),
      ("Bon hochladen","On-Device-OCR — privat, keine Quota, offline"),
      ("Dein Vorrat","was im Haus ist, Ablauf-bewusst — der Startbildschirm"),
      ("Tages-Log","ein Tap, um Gegessenes abzuhaken"),
      ("Insights","Wochen-Lücken + „nutze, was du hast“ + Next Cart"),
      ("Nutri-Coach + Löschen","Begründung in Klartext · DSGVO-Löschung jederzeit")]
x0=0.9; yy=2.05
for i,(t,sub) in enumerate(flow):
    y=yy+i*0.73
    rect(s,x0,y,0.5,0.5,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,x0,y-0.02,0.5,0.5,[[R(str(i+1),15,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x0+0.7,y-0.04,5.7,0.4,[[R(t,15,INK,True,False),R("   —  "+sub,11.5,INK_SOFT,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if i<5: rect(s,x0+0.23,y+0.5,0.03,0.23,fill=SAGE_LT)
rx=7.25
rect(s,rx,1.95,5.2,4.55,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
rect(s,rx,1.95,5.2,0.66,fill=SAGE,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.05)
txt(s,rx+0.3,1.99,4.6,0.6,[[R("NutriWise  ·  Deine Ergebnisse",13,WHITE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
txt(s,rx+0.3,2.78,4.6,0.5,[[R("Schätzung aus Einkauf + Tages-Log — keine medizinische Beratung",10.5,INK_SOFT,False,True)]])
rect(s,rx+0.3,3.25,4.6,1.5,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08)
txt(s,rx+0.55,3.4,4.1,0.4,[[R("DEIN NEXT CART",11,SAGE_DK,True,False)]])
txt(s,rx+0.55,3.75,4.1,0.5,[[R("+ Rote Linsen (getrocknet)",17,INK,True,False)]])
txt(s,rx+0.55,4.2,4.1,0.5,[[R("Schließt deine Ballaststoff- & Pflanzenprotein-Lücke — und ist erdnussfrei.",11.5,INK_SOFT,False,False)]],ls=1.05)
for i,(lab,val) in enumerate([("Ballastst.","niedrig"),("Protein","ok"),("Zucker","ok")]):
    bx=rx+0.3+i*1.55
    rect(s,bx,4.95,1.4,0.85,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.12)
    txt(s,bx,5.05,1.4,0.4,[[R(lab,11,INK_SOFT,True,False)]],align=PP_ALIGN.CENTER)
    txt(s,bx,5.35,1.4,0.4,[[R(val,15,CLAY if val=="niedrig" else SAGE_DK,True,False)]],align=PP_ALIGN.CENTER)
txt(s,rx,6.05,5.2,0.4,[[R("illustrativ — live in der Demo gezeigt",10,INK_SOFT,False,True)]],align=PP_ALIGN.CENTER)
footer(s,7)
notes(s,"Sprecher 1 → übergibt an Sprecher 2",
      "LIVE-DEMO (~2 Min). Zuerst den Positioning-Satz, dann: Onboarding-Chat (Profil) → Bon hochladen "
      "(On-Device-OCR — betonen: privat, keine API-Quota, offline) → auf dem Vorrat landen (das Home) → "
      "Gegessenes per Tap loggen → Insights: Wochen-Lücken + „nutze, was du hast“ + Next Cart → "
      "Nutri-Coach-Formulierung + DSGVO-Löschung. HINWEIS: Bons hängen nicht mehr an Gemini, dort kein "
      "Quota-Risiko. Nur der Nutri-Coach nutzt Gemini (nur Formulierung) — bei erschöpftem Tageslimit "
      "funktionieren Zahlen & Empfehlung weiter, nur der Text degradiert. Dann Übergabe an Sprecher 2.")

# ================= FOLIE 8 — Evaluations-Metrik
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Evaluations-Metrik"); heading(s,0.9,0.92,11.7,"Das Schwierige: kryptischen Bon-Text zu echtem Lebensmittel machen",size=25); rule(s,0.92,1.62)
txt(s,0.9,1.95,6.2,3.7,
    [[R("„AB Bud 6x0,3l MW“",17,CLAY,True,False),R("   →   ?",17,INK,True,False)],
     [R("",6,INK,False,False)],
     [R("Jede Bon-Zeile muss ein echtes Lebensmittel mit echten Nährwerten werden. Es gibt ",14.5,INK,False,False),
      R("kein Ground-Truth-Label",14.5,INK,True,False),R(" für „welches Lebensmittel ist das wirklich?“",14.5,INK,False,False)],
     [R("",5,INK,False,False)],
     [R("Das ist regelbasiert / LLM-Hybrid, ",14.5,INK,False,False),
      R("kein trainiertes ML-Modell",14.5,SAGE_DK,True,False),
      R(" — also bauten wir eine eigene Evaluation: ein ",14.5,INK,False,False),
      R("LLM-Judge-Panel.",14.5,SAGE_DK,True,False)],
     [R("3 unabhängige Judges bewerten jeden Match — ",14.5,INK,False,False),
      R("korrekt / teilweise / falsch / kein Treffer",14.5,INK,True,False),
      R(" — die Mehrheit gewinnt.",14.5,INK,False,False)]],ls=1.12)
px=7.5; rect(s,px,1.95,4.95,3.85,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05)
txt(s,px+0.35,2.17,4.3,0.4,[[R("WAS WIR MESSEN",12,SAGE_DK,True,False)]])
stat(s,px+0.3,2.62,4.35,1.32,"Precision","Match-Qualität","% als korrekt bewertet — oder nutzbar (korrekt + teilw.)",accent=SAGE)
stat(s,px+0.3,4.1,4.35,1.32,"Konsens","Verlässlichkeit","% der Items, bei denen alle 3 Judges übereinstimmen",accent=SAGE_DK)
footer(s,8)
notes(s,"Sprecher 2",
      "Sprecher 2 (Jennifer) übernimmt. Klar sagen: das ist KEIN klassisches Supervised-ML-Projekt — es ist "
      "regelbasiert / LLM-Hybrid. Kernaufgabe: kryptischen Bon-Text wie „AB Bud 6x0,3l MW“ auf ein echtes "
      "Lebensmittel matchen. Es gibt keine Ground Truth, also ist unsere Evaluation ein LLM-Judge-Panel: 3 "
      "unabhängige Judges, Mehrheitsentscheid. Zwei Zahlen: Precision (korrekt/nutzbar) und Judge-Konsens "
      "(vertrauen wir der Metrik selbst). Produktziel: ≥70% der Items nutzbar.")

# ================= FOLIE 9 — Daten
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Daten — Quellen, Form & Challenges"); heading(s,0.9,0.92,11.5,"Drei Datenwelten, zusammengefügt"); rule(s,0.92,1.62)
srcs=[("Kassenbons","Echte deutsche Bons (Lidl / Netto / Edeka) → On-Device-OCR (Tesseract) → Items. Privat: keine Bondaten verlassen das Gerät."),
      ("Open Food Facts","Crowdsourced-Marken-API + lokaler Cache. Erste Identitäts-Zuordnung; pro Feld lückenhaft."),
      ("BLS-4.0-Tabelle","Deutsche Bundes-Nährwert-DB. ~7.140 Zeilen × 418 Spalten — 7 Makros + 11 Mikros extrahiert.")]
cw=3.63;x0=0.9;yy=1.95
for i,(t,b) in enumerate(srcs):
    x=x0+i*(cw+0.24)
    rect(s,x,yy,cw,1.9,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    rect(s,x,yy+0.14,0.08,1.62,fill=SAGE if i<2 else SAGE_DK)
    txt(s,x+0.32,yy+0.2,cw-0.5,0.4,[[R(t,15,INK,True,False)]])
    txt(s,x+0.32,yy+0.72,cw-0.55,1.1,[[R(b,11.5,INK_SOFT,False,False)]],ls=1.1)
rect(s,0.9,4.1,5.55,1.7,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,1.15,4.25,5.1,0.4,[[R("EVALUATIONS-SET (Form)",11.5,SAGE_DK,True,False)]])
txt(s,1.15,4.67,5.15,1.1,
    [[R("76",20,SAGE_DK,True,False),R(" einzigartige Items   ·   ",12.5,INK,False,False),
      R("227",20,SAGE_DK,True,False),R(" Zeilen-Vorkommen",12.5,INK,False,False)],
     [R("13 Spalten: Rohname & normalisiert, Kategorie, Vorkommen,",11,INK_SOFT,False,False)],
     [R("off/bls Match-Name · Verdikt · Konsens · n_judges",11,INK_SOFT,False,False)]],ls=1.12,space_after=3)
rect(s,6.7,4.1,5.73,1.7,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05)
txt(s,6.95,4.25,5.3,0.4,[[R("EHRLICHE DATEN-CHALLENGES",11.5,CLAY,True,False)]])
for i,t in enumerate(["Kryptische OCR / Marke-vs-generisch („Budweiser“ → „Pilsner Bier“)",
                      "OFF-Zucker = Gesamt-, nicht freier Zucker; BLS ohne NOVA",
                      "Dünne Nährstoff-Coverage → wenige Nährstoffe gut, nicht viele halbgar"]):
    txt(s,6.95,4.62+i*0.37,5.35,0.4,[[R("•  ",12,CLAY,True,False),R(t,11.5,INK,False,False)]],ls=1.0)
footer(s,9)
notes(s,"Sprecher 2",
      "Drei Datenwelten: (1) echte deutsche Bons → On-Device-OCR (Tesseract) → Items — privat, keine API, "
      "keine Quota, (2) Open Food Facts als erste Identitäts-Zuordnung, (3) die deutsche Bundes-BLS-4.0-"
      "Tabelle (7.140 Zeilen, wir ziehen 7 Makros + 11 Mikros) für die Nährwerte. Eval-Set: 76 Items aus "
      "227 echten Vorkommen, 13 Spalten. Die Challenges SIND die Story — kryptische OCR, Marke-vs-generisch, "
      "OFF-Zucker ist Gesamt- statt freier Zucker, BLS ohne Processing-Rating, dünne Coverage.")

# ================= FOLIE 10 — EDA
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Explorative Datenanalyse"); heading(s,0.9,0.92,11.5,"Wie echte Warenkörbe wirklich aussehen"); rule(s,0.92,1.62)
rect(s,0.9,1.95,5.75,4.3,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
s.shapes.add_picture(CH("eda_categories_de.png"),Inches(1.05),Inches(2.12),height=Inches(3.9))
rect(s,6.85,1.95,5.58,4.3,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
s.shapes.add_picture(CH("eda_tail_de.png"),Inches(7.0),Inches(2.12),height=Inches(3.9))
txt(s,0.9,6.32,11.5,0.5,
    [[R("Erkenntnis:  ",12.5,SAGE_DK,True,False),
      R("Warenkörbe sind obst-/gemüse- und milchlastig, aber der Long-Tail ist brutal — die Hälfte der "
        "Items erscheint nur einmal, der Matcher muss also generalisieren, nicht auswendig lernen.",12.5,INK,False,False)]],ls=1.05)
footer(s,10)
notes(s,"Sprecher 2",
      "Ehrlich anmerken: diese Plots gab es im Repo nicht — wir haben sie für dieses Deck aus den "
      "Evaluationsdaten gebaut. Zwei Befunde: links Käufe stark bei Obst/Gemüse & Milch (die Basics); "
      "rechts ein Long-Tail — 50% der Items erscheinen nur einmal. Wir können also nicht auswendig lernen; "
      "der Matcher muss auf einmalige, kryptisch abgekürzte Produkte generalisieren. Das prägte die Pipeline.")

# ================= FOLIE 11 — Pipeline
s=slide(); contentbg(s); kicker(s,0.9,0.6,"ML-Workflow / Daten-Pipeline"); heading(s,0.9,0.92,11.7,"Ein gestufter Resolver — günstig & sicher zuerst, Fallback zuletzt"); rule(s,0.92,1.62)
stages=[("Bon","Bild / Text"),("Extraktion","On-Device-OCR → Items"),("Normalisieren","Menge · Preis · Non-Food raus")]
tiers=[("Tier 0","Gelernte verifizierte Matches",SAGE_DK,"menschlich bestätigt · Konf. 1.0"),
       ("Tier 1","OpenFoodFacts + Fuzzy-Match",SAGE,"→ BLS-Bridge für DE-Nährwerte"),
       ("Tier 1b","BLS-Grundlebensmittel (Rohware)",SAGE_LT,"per Typ-Übereinstimmung abgesichert"),
       ("Tier 3","Kategorie-Fallback",SAND,"niedrige Konfidenz · „unbekannt“ markiert")]
yy=2.05;bw=2.55;x0=0.9
for i,(t,sub) in enumerate(stages):
    x=x0+i*(bw+0.55)
    rect(s,x,yy,bw,1.0,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08)
    txt(s,x,yy+0.18,bw,0.4,[[R(t,15,WHITE,True,False)]],align=PP_ALIGN.CENTER)
    txt(s,x,yy+0.6,bw,0.35,[[R(sub,11,SAGE_LT,False,False)]],align=PP_ALIGN.CENTER)
    if i<2: txt(s,x+bw+0.04,yy+0.2,0.5,0.6,[[R("→",24,SAGE,True,False)]])
txt(s,x0+3*(bw+0.55)-0.35,yy+0.2,1.4,0.6,[[R("→",24,SAGE,True,False)]])
ty=3.35
for i,(tag,name,col,sub) in enumerate(tiers):
    y=ty+i*0.72
    rect(s,0.9,y,8.4,0.62,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.12)
    rect(s,0.9,y,1.55,0.62,fill=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.12)
    tc=INK if col in (SAGE_LT,SAND) else WHITE
    txt(s,0.9,y,1.55,0.62,[[R(tag,13,tc,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,2.65,y,3.6,0.62,[[R(name,13.5,INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,6.0,y,3.2,0.62,[[R(sub,11.5,INK_SOFT,False,True)]],anchor=MSO_ANCHOR.MIDDLE)
rect(s,9.65,3.35,2.78,2.78,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,9.85,3.55,2.4,0.4,[[R("AUSGABE",11.5,SAGE_DK,True,False)]])
for i,t in enumerate(["Nährwerte (Makros + Mikros)","Konfidenz-Label","Quellenangabe","→ Lücken + Next Cart"]):
    txt(s,9.85,4.0+i*0.53,2.5,0.5,[[R("• ",12,SAGE,True,False),R(t,12,INK,False,False)]],ls=1.0)
txt(s,0.9,6.35,8.4,0.4,[[R("Bewusst regelbasiert — kein gefittetes Modell — damit die Logik transparent & vertrauenswürdig bleibt.",12,INK_SOFT,False,True)]])
footer(s,11)
notes(s,"Sprecher 2",
      "Wegen des Long-Tails nutzen wir kein einzelnes Modell — sondern einen gestuften Resolver. Bon → "
      "On-Device-OCR (Tesseract, ersetzt Gemini — privat, keine Quota) → Normalisierung. Dann vier Tiers, "
      "günstigste/sicherste zuerst: Tier 0 menschlich verifiziert, Tier 1 OFF-Fuzzy-Match mit Bridge zu "
      "deutschen BLS-Nährwerten, Tier 1b Rohware, Tier 3 ehrlicher Low-Confidence-Kategorie-Fallback. Jede "
      "Ausgabe trägt Nährwerte + Konfidenz + Quelle. Klar sagen: bewusst regelbasiert — Regeln sind "
      "transparenter und vertrauenswürdiger als eine Blackbox für ein Gesundheits-MVP. Kein gefittetes "
      "Modell im Matching-Pfad; die einzige LLM (Gemini) ist der Nutri-Coach, der nur formuliert.")

# ================= FOLIE 12 — Baseline
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Baseline-Modell & Performance"); heading(s,0.9,0.92,11.7,"Zwei Matching-Strategien, direkt verglichen"); rule(s,0.92,1.62)
rect(s,0.9,1.9,7.55,3.55,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
s.shapes.add_picture(CH("baseline_de.png"),Inches(1.05),Inches(2.1),width=Inches(7.25))
stat(s,8.75,1.9,3.68,1.1,"66%","OFF nutzbar","korrekt + teilweise — unsere Baseline",accent=SAGE)
stat(s,8.75,3.12,3.68,1.1,"2×","OFF schlägt BLS","54% vs 32% korrekt bei Identität",accent=SAGE_DK)
stat(s,8.75,4.34,3.68,1.1,"96%","Judge-Konsens","3/3 einstimmig — Metrik ist verlässlich",accent=SAGE)
rect(s,0.9,5.62,11.53,0.9,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,1.2,5.72,11,0.72,
    [[R("Warum die Architektur:  ",13,SAGE_DK,True,False),
      R("OFF gewinnt bei der Identität und führt. BLS matcht schlechter (~⅓ falsches Lebensmittel), daher "
        "werden seine Nährwerte nur unter Guard geliehen — nie als primärer Matcher. 66% nutzbar liegt knapp unter der 70%-Marke.",
        13,INK,False,False)]],ls=1.08,anchor=MSO_ANCHOR.MIDDLE)
footer(s,12)
notes(s,"Sprecher 2",
      "Wichtige Klarstellung: das sind zwei UNABHÄNGIGE Matching-Strategien im direkten Vergleich, KEINE "
      "kumulative Verbesserung. Vom Panel bewertet: OpenFoodFacts-Identitäts-Match = 54% korrekt / 66% "
      "nutzbar; die BLS-Tabelle allein = nur 32% korrekt, 36% falsch. Judge-Konsens 96% einstimmig, also "
      "vertrauen wir der Zahl. Genau dieser Vergleich ist der Grund, warum die Architektur OFF für die "
      "Identität nutzt und BLS-Nährwerte nur unter Typ-Übereinstimmungs-Guard leiht — eine lose Bridge "
      "würde BLS' ~33%-Falschlebensmittel-Problem zurückbringen. 66% nutzbar liegt knapp unter der "
      "70%-Marke — genau diese Lücke schließt der gelernte Tier-0-Store.")

# ================= FOLIE 13 — Nächstes / Zukunft
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Nächste Woche & darüber hinaus"); heading(s,0.9,0.92,11.5,"Von der Baseline zur Vertrauenswürdigkeit"); rule(s,0.92,1.62)
cols=[("Besseres Matching",SAGE,["Fuzzy-Baseline 66% → Embedding-Re-Ranker, Ziel >70% nutzbar",
        "Gelernter Tier-0-Store wird mit jeder Korrektur besser","Personalisierte Portions- & Verbrauchsschätzung (ML)"]),
      ("Validieren & wachsen",SAGE_DK,["Live-User-Tests — klebt die Abhak-Gewohnheit?",
        "A/B-Test: bester Moment für „was gegessen?“","10-KPI-Funnel instrumentieren (Upload → Insight → Retention)"])]
cw=5.75;x0=0.9;yy=2.05
for i,(t,ac,items) in enumerate(cols):
    x=x0+i*(cw+0.13)
    rect(s,x,yy,cw,3.75,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
    rect(s,x,yy,cw,0.72,fill=ac,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.06)
    txt(s,x+0.3,yy+0.08,cw-0.5,0.6,[[R(t,17,WHITE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    for j,it in enumerate(items):
        yj=yy+1.05+j*0.92
        rect(s,x+0.35,yj+0.03,0.3,0.3,fill=ac,shape=MSO_SHAPE.OVAL)
        txt(s,x+0.35,yj+0.0,0.3,0.3,[[R("→",12,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+0.85,yj-0.05,cw-1.2,0.9,[[R(it,13.5,INK,False,False)]],ls=1.05)
txt(s,0.9,6.15,11.5,0.4,[[R("Prinzip:  ",12,SAGE_DK,True,False),
    R("ML dort, wo Regeln nicht generalisieren (Matching, Portionen) — Lücken & Empfehlungen bleiben regelbasiert und transparent.",12,INK_SOFT,False,True)]])
footer(s,13)
notes(s,"Sprecher 2",
      "Zwei Tracks nächste Woche. Matching: den Resolver zum lernenden System machen — jede Nutzerkorrektur "
      "füttert den Tier-0-Store — plus bessere Normalisierung, um über 70% nutzbar zu kommen. Produkt: "
      "Live-User-Tests starten und den KPI-Funnel instrumentieren, damit wir die Trust-Hypothese mit echtem "
      "Verhalten validieren, nicht nur mit einer steigenden Genauigkeitszahl.")

# ================= FOLIE 14 — Danke
s=slide(); bg(s,CANVAS)
rect(s,0,0,13.333,0.14,fill=SAGE)
rect(s,0,7.36,13.333,0.14,fill=SAGE)
logo(s,0.9,0.9)
txt(s,0.9,2.5,11.5,1.2,[[R("Danke.",50,INK,True,False)]])
rule(s,0.92,3.7,3.4)
txt(s,0.92,3.95,11,0.6,[[R("Lade deinen Kassenbon hoch. Bekomm eine klügere Sache für den nächsten Einkauf.",18,SAGE_DK,False,True)]])
txt(s,0.9,4.95,11.5,1.0,
    [[R("Jennifer Rake",15,INK,True,False),R("   &   ",15,INK_SOFT,False,False),R("Stuart Kasemeier",15,INK,True,False)],
     [R("NutriWise  ·  AI PM Capstone  ·  Midterm V1  —  Fragen & Feedback willkommen",12.5,INK_SOFT,False,False)]],space_after=8)
notes(s,"Beide",
      "Mit dem Einzeiler schließen und Fragen einladen. Bei Nachfragen zur Genauigkeit: ehrliche Baseline "
      "(OFF 66% nutzbar), klarer Plan (gelernter Tier-0) zum Schließen der Lücke — und die schwierigere "
      "Validierungsfrage ist Vertrauen/Verhalten, die als Nächstes die Live-User-Tests beantworten.")

out=os.path.join(OUTDIR,"NutriWise_Midterm_Presentation_v1_DE.pptx")
prs.save(out); print("saved",out,"slides:",len(prs.slides._sldIdLst))
