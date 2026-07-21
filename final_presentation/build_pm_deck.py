#!/usr/bin/env python3
"""NutriWise — PM Meeting 3 status deck (ENGLISH). Corporate design (sage, sans-serif).
8 slides, no charts. Content = PM Meeting slides, reconciled with current repo facts
(on-device OCR already replaced Gemini; baseline 76 items / 227 occurrences)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------- corporate palette (shared with build_deck*.py)
SAGE     = "7C9A6A"; SAGE_DK  = "5F7B4F"; SAGE_LT  = "B7C9A9"; SAGE_SOFT= "EEF2EA"
CLAY     = "B36A4A"; SAND     = "CFC9BA"; INK      = "1D1D21"; INK_SOFT = "6E6F74"
CANVAS   = "FAFAFA"; SURFACE  = "FFFFFF"; LINE     = "E5E4E0"; WHITE    = "FFFFFF"
FONT     = "Arial"
NSLIDES  = 9

def C(h): return RGBColor.from_string(h)
HERE   = os.path.dirname(os.path.abspath(__file__))
OUTDIR = HERE

# ---------------------------------------------------- helpers (ported from build_deck_de.py)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def bg(s,h): s.background.fill.solid(); s.background.fill.fore_color.rgb=C(h)
def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,shadow=False,adj=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=C(line); sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if adj is not None:
        try: sp.adjustments[0]=adj
        except Exception: pass
    if shadow:
        from pptx.oxml.ns import qn
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{})
        sh=el.makeelement(qn('a:outerShdw'),{'blurRad':'80000','dist':'30000','dir':'5400000','rotWithShape':'0'})
        clr=el.makeelement(qn('a:srgbClr'),{'val':'1D1D21'}); alp=el.makeelement(qn('a:alpha'),{'val':'12000'})
        clr.append(alp); sh.append(clr); ef.append(sh); el.append(ef)
    return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space_after=4,ls=1.0,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,col,b,it) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it
            r.font.color.rgb=C(col); r.font.name=FONT
    return tb
def R(t,sz,col=INK,b=False,it=False): return (t,sz,col,b,it)
def notes(s,speaker,body): s.notes_slide.notes_text_frame.text=f"[{speaker}]\n{body}"
def kicker(s,x,y,t): txt(s,x,y,10,0.3,[[R(t.upper(),12.5,SAGE_DK,True,False)]])
def heading(s,x,y,w,t,size=29): txt(s,x,y,w,1.0,[[R(t,size,INK,True,False)]],ls=1.0)
def rule(s,x,y,w=1.5): rect(s,x,y,w,0.045,fill=SAGE)
def footer(s,page):
    txt(s,0.6,7.03,5,0.3,[[R("NutriWise",9,INK_SOFT,True,False),R("  ·  PM Meeting 3",9,INK_SOFT,False,False)]])
    txt(s,10.8,7.03,1.9,0.3,[[R(f"{page:02d} / {NSLIDES:02d}",9,INK_SOFT,False,False)]],align=PP_ALIGN.RIGHT)
def logo(s,x,y,dark=False):
    rect(s,x,y,0.34,0.34,fill=SAGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.3)
    txt(s,x+0.46,y-0.03,4,0.42,[[R("NUTRIWISE",13,WHITE if dark else INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
def contentbg(s): bg(s,CANVAS)
def stat(s,x,y,w,h,big,label,sub=None,accent=SAGE):
    rect(s,x,y,w,h,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.06)
    rect(s,x,y+0.14,0.08,h-0.28,fill=accent)
    body=[[R(big,25,accent,True,False)],[R(label,12.5,INK,True,False)]]
    if sub: body.append([R(sub,10.5,INK_SOFT,False,False)])
    txt(s,x+0.3,y+0.14,w-0.42,h-0.26,body,anchor=MSO_ANCHOR.MIDDLE,space_after=3,ls=1.02)
def pill(s,x,y,w,t,fill=SAGE,col=WHITE):
    rect(s,x,y,w,0.44,fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    txt(s,x,y+0.02,w,0.4,[[R(t,11,col,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ================= SLIDE 1 — Title
s=slide(); bg(s,CANVAS)
rect(s,0,0,13.333,0.14,fill=SAGE)
rect(s,8.9,0.14,4.433,7.36,fill=SAGE_SOFT)
logo(s,0.9,0.85)
txt(s,0.9,2.05,7.7,2.4,
    [[R("Supermarket Optimizer",37,INK,True,False)],
     [R("Project status &",37,SAGE_DK,True,False)],
     [R("priorities.",37,SAGE_DK,True,False)]],ls=1.05)
rule(s,0.92,4.45,3.4)
txt(s,0.92,4.7,7.5,0.7,[[R("Where we are, what we've solved, and what ships next — a receipt-based nutrition assistant.",15,INK_SOFT,False,True)]])
txt(s,0.9,5.65,7.8,1.2,
    [[R("Jennifer Rake",15,INK,True,False),R("  ·  UI / MVP      ",13,SAGE_DK,True,False),
      R("Stuart Kasemeier",15,INK,True,False),R("  ·  Backend",13,INK_SOFT,True,False)],
     [R("AI Product Management Capstone  ·  PM Meeting 3  ·  July 21, 2026",12.5,INK_SOFT,False,False)]],space_after=8)
# agenda rail (right panel)
txt(s,9.5,1.95,3.2,0.4,[[R("TODAY'S AGENDA",11.5,SAGE_DK,True,False)]])
for i,t in enumerate(["Business Idea","Project Timeline","Product Breakdown","Challenges","Team Collaboration","Priorities"]):
    yy=2.5+i*0.72
    rect(s,9.5,yy,0.26,0.26,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,9.95,yy-0.06,2.9,0.4,[[R(t,13,INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if i<5: rect(s,9.62,yy+0.3,0.02,0.42,fill=SAGE_LT)
notes(s,"Both — Stuart opens",
      "Open the PM meeting: quick status update on Supermarket Optimizer / NutriWise. Six blocks: business "
      "idea, timeline, product breakdown, challenges (solved + upcoming), team collaboration, and priorities "
      "for this and next week. Main headline for today: we want to get to live testing on our phones.")

# ================= SLIDE 2 — Business Idea
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Business Idea"); heading(s,0.9,0.92,11.5,"Better nutrition with minimal user effort"); rule(s,0.92,1.62)
steps=[("1","Upload a receipt","Snap or drop a grocery receipt — no manual food diary."),
       ("2","AI analyzes","Products and their nutrition are extracted automatically."),
       ("3","Find the gaps","We spot the nutrients the week is missing."),
       ("4","Recommend next","Healthier picks for the next shopping trip, with a reason.")]
cw=2.86;gp=0.24;x0=0.9;yy=2.05
for i,(n,t,body) in enumerate(steps):
    x=x0+i*(cw+gp)
    rect(s,x,yy,cw,2.85,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    rect(s,x,yy,cw,0.7,fill=SAGE if i<3 else SAGE_DK,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.08)
    txt(s,x,yy+0.06,cw,0.6,[[R("STEP "+n,12,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.25,yy+0.88,cw-0.5,0.5,[[R(t,16,INK,True,False)]])
    txt(s,x+0.25,yy+1.5,cw-0.5,1.3,[[R(body,12.5,INK_SOFT,False,False)]],ls=1.08)
    if i<3: txt(s,x+cw-0.02,yy+1.05,0.34,0.5,[[R("›",26,SAGE,True,False)]])
rect(s,0.9,5.25,11.53,1.0,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08)
txt(s,1.2,5.4,11,0.72,
    [[R("Why the receipt?  ",13.5,SAGE_DK,True,False),
      R("It's a passive, honest record of what people actually buy — so users get personalized nutrition "
        "guidance without the manual logging that kills every other diet app.",13.5,INK,False,False)]],
    ls=1.12,anchor=MSO_ANCHOR.MIDDLE)
footer(s,2)
notes(s,"Stuart",
      "The core idea in one line: upload a receipt, get smarter nutrition for the next trip. The wedge is "
      "effort — the receipt is a passive proxy for what was really eaten, so no manual food logging. Goal: "
      "better nutrition with minimal user effort.")

# ================= SLIDE 3 — Project Timeline
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Project Timeline"); heading(s,0.9,0.92,11.5,"Three phases, one shipping MVP"); rule(s,0.92,1.62)
ty=2.5; rect(s,1.1,ty+0.98,11.2,0.05,fill=LINE)
cols=[("WEEK 1",SAGE_DK,"MVP trials",["Checked Claude's capabilities","Ran our first local MVP","Identified shortcomings","Refined documentation"]),
      ("WEEK 2+3",SAGE,"In-depth work",["Clear, structured documentation","Structured rework of the current MVP","Switching partners per Epic again"]),
      ("WEEK 3",INK_SOFT,"Testing & polish",["Deployment","Live testing on real devices","GDPR compliance"])]
cw=3.63;x0=1.1
for i,(tag,col,when,items) in enumerate(cols):
    x=x0+i*(cw+0.24)
    rect(s,x+cw/2-0.13,ty+0.86,0.3,0.3,fill=col,shape=MSO_SHAPE.OVAL)
    pill(s,x,ty-0.05,cw*0.5,tag,fill=col)
    txt(s,x+cw*0.53,ty-0.02,cw*0.5,0.44,[[R(when,12.5,INK_SOFT,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    cy=ty+1.5
    rect(s,x,cy,cw,2.9,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    for j,it in enumerate(items):
        yy=cy+0.28+j*0.58
        mark="✓" if i==0 else ("▸" if i==1 else "•")
        txt(s,x+0.3,yy,0.4,0.4,[[R(mark,13,col,True,False)]])
        txt(s,x+0.7,yy-0.02,cw-0.95,0.6,[[R(it,12.5,INK,False,False)]],ls=1.0)
footer(s,3)
notes(s,"Stuart",
      "Three-phase plan. Week 1 was MVP trials — probing Claude, a first local build, spotting gaps, "
      "tightening docs. Weeks 2–3 are the deeper pass: clean documentation and a structured rework of the "
      "MVP, back to swapping partners per Epic. Week 3 closes with deployment, live testing, and GDPR.")

# ================= SLIDE 4 — Product Breakdown
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Product Breakdown"); heading(s,0.9,0.92,11.5,"Six building blocks, end to end"); rule(s,0.92,1.62)
blocks=[("Login / Auth","Secure sign-in and account handling."),
        ("Onboarding","Chat-style profile generation — goal, needs, exclusions."),
        ("Receipt Upload","PDF, image, and (soon) camera capture."),
        ("Extraction & Matching","On-device OCR → items → nutrition (rule/LLM hybrid)."),
        ("Dashboard & Recs","Nutrient gaps + healthier next-trip recommendations."),
        ("Recipe Generation","Turns improvements into concrete, intuitive meals.")]
cw=3.72;ch=1.9;gx=0.28;gy=0.34;x0=0.9;y0=2.0
for i,(t,b) in enumerate(blocks):
    c=i%3; r=i//3
    x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    rect(s,x,y,cw,ch,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.06)
    rect(s,x+0.28,y+0.26,0.5,0.5,fill=SAGE if i<5 else SAGE_DK,shape=MSO_SHAPE.OVAL)
    txt(s,x+0.28,y+0.24,0.5,0.5,[[R(str(i+1),15,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.95,y+0.28,cw-1.1,0.5,[[R(t,15.5,INK,True,False)]])
    txt(s,x+0.95,y+0.82,cw-1.15,0.95,[[R(b,12,INK_SOFT,False,False)]],ls=1.08)
footer(s,4)
notes(s,"Stuart / Jenny",
      "The product is six blocks: login/auth, chat-style onboarding that builds the profile, receipt upload "
      "(PDF and image today, camera coming), extraction & matching, the dashboard with gaps and "
      "recommendations, and recipe generation to make the improvements intuitive. Extraction now runs "
      "fully on-device.")

# ================= SLIDE 5 — Challenges & Upcoming
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Challenges"); heading(s,0.9,0.92,11.7,"What we hit, how we solved it, what's next"); rule(s,0.92,1.62)
# left: encountered + solutions
lx=0.9; lw=5.55
rect(s,lx,1.95,lw,4.35,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
txt(s,lx+0.35,2.15,lw-0.6,0.4,[[R("ENCOUNTERED",11.5,CLAY,True,False)]])
for i,t in enumerate(["Product matching was the hard part","Too little matching logic up front"]):
    txt(s,lx+0.35,2.55+i*0.42,lw-0.7,0.4,[[R("•  ",12,CLAY,True,False),R(t,13,INK,False,False)]],ls=1.0)
rect(s,lx+0.35,3.5,lw-0.7,0.02,fill=LINE)
txt(s,lx+0.35,3.62,lw-0.6,0.4,[[R("SOLUTIONS",11.5,SAGE_DK,True,False)]])
for i,t in enumerate(["Enabled manual matching — collect data now, train a smart matcher that needs no interaction later",
                      "Rewrote the documentation and reworked the MVP"]):
    yy=4.02+i*0.9
    rect(s,lx+0.35,yy+0.02,0.32,0.32,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,lx+0.35,yy-0.02,0.32,0.32,[[R("✓",12,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+0.82,yy-0.03,lw-1.2,0.85,[[R(t,12.5,INK,False,False)]],ls=1.06)
# right: upcoming
rx=6.75; rw=5.68
rect(s,rx,1.95,rw,4.35,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.04)
txt(s,rx+0.35,2.15,rw-0.6,0.4,[[R("UPCOMING CHALLENGES",11.5,SAGE_DK,True,False)]])
ups=[("Deployment","Ship a stable hosted build."),
     ("GDPR compliance","Export / delete, privacy by design."),
     ("Extraction robustness","On-device OCR already replaced Gemini — next: make it reliable across every supermarket format."),
     ("Camera capture","Take a photo in-app, alongside PDF & image upload.")]
for i,(t,b) in enumerate(ups):
    yy=2.62+i*0.9
    rect(s,rx+0.35,yy+0.03,0.32,0.32,fill=SAGE_DK,shape=MSO_SHAPE.OVAL)
    txt(s,rx+0.35,yy-0.01,0.32,0.32,[[R("▸",12,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,rx+0.82,yy-0.03,rw-1.2,0.85,[[R(t+"  —  ",12.5,INK,True,False),R(b,12,INK_SOFT,False,False)]],ls=1.05)
footer(s,5)
notes(s,"Stuart / Jenny",
      "The big one we hit was product matching — we didn't put enough logic in up front. Fix: enable manual "
      "matching now so we collect labelled data, then train a smarter matcher that needs no user interaction. "
      "We also rewrote docs and reworked the MVP. Upcoming: deployment, GDPR, and extraction robustness — "
      "important reconciliation: we've ALREADY moved extraction on-device (Tesseract/PyMuPDF), so Gemini is "
      "out of the extraction path; the remaining work is generalizing across supermarket formats. Plus "
      "in-app camera capture next to PDF/image upload.")

# ================= SLIDE 6 — Current focus / open technical items
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Current focus — open items"); heading(s,0.9,0.92,11.7,"What we're actively working on right now"); rule(s,0.92,1.62)
items=[("Supabase load time","App data loads too slowly.","Caching, leaner queries & lazy loading to cut wait time."),
       ("Recipes & recommendations","Generation quality and speed.","Sharper prompts + ranking so picks are faster and more relevant."),
       ("Portion & quantity tracking","Pantry counts pack sizes (e.g. 3× milk).","Track quantities so stock and gaps stay accurate over the week."),
       ("Image→text & DB matching","Reliably turn a receipt into nutrients.","Train & tune on-device OCR and product→nutrition matching.")]
cw=5.75;ch=1.9;gx=0.13;gy=0.28;x0=0.9;y0=2.0
for i,(t,prob,sol) in enumerate(items):
    c=i%2; r=i//2
    x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    rect(s,x,y,cw,ch,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.05)
    rect(s,x,y+0.14,0.08,ch-0.28,fill=SAGE_DK if i==0 else SAGE)
    rect(s,x+0.3,y+0.28,0.5,0.5,fill=SAGE_SOFT,shape=MSO_SHAPE.OVAL)
    txt(s,x+0.3,y+0.26,0.5,0.5,[[R(str(i+1),16,SAGE_DK,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.95,y+0.24,cw-1.15,0.4,[[R(t,15.5,INK,True,False)]])
    txt(s,x+0.95,y+0.68,cw-1.15,0.4,[[R(prob,12,INK_SOFT,False,True)]],ls=1.02)
    txt(s,x+0.95,y+1.12,cw-1.15,0.7,[[R("→  ",12.5,SAGE_DK,True,False),R(sol,12.5,INK,False,False)]],ls=1.05)
footer(s,6)
notes(s,"Both",
      "The concrete open items we're heads-down on right now, with how we're tackling each: (1) Supabase "
      "data load time in the app — caching, leaner queries, lazy loading; (2) recipe & recommendation "
      "generation — quality and speed via better prompts and ranking; (3) portion / quantity tracking in the "
      "pantry, e.g. counting 3× milk, so stock and gaps stay accurate; (4) training and tuning image→text "
      "(on-device OCR) plus database matching to reliably pull ingredients / nutrients. These feed directly "
      "into the live-testing goal.")

# ================= SLIDE 7 — Team Collaboration
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Team Collaboration"); heading(s,0.9,0.92,11.5,"How we work together"); rule(s,0.92,1.62)
# task distribution
lx=0.9; lw=5.75; yy=2.05
rect(s,lx,yy,lw,4.2,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
rect(s,lx,yy,lw,0.7,fill=SAGE,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.06)
txt(s,lx+0.3,yy+0.06,lw-0.5,0.6,[[R("TASK DISTRIBUTION",13,WHITE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
people=[("Jenny","UI expert · MVP refinement"),("Stuart","Backend expert · documentation")]
for i,(nm,role) in enumerate(people):
    py=yy+1.05+i*1.0
    rect(s,lx+0.35,py,0.6,0.6,fill=SAGE_SOFT,shape=MSO_SHAPE.OVAL)
    txt(s,lx+0.35,py,0.6,0.6,[[R(nm[0],17,SAGE_DK,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+1.15,py+0.02,lw-1.4,0.6,[[R(nm,15,INK,True,False)],[R(role,12,INK_SOFT,False,False)]],space_after=1)
rect(s,lx+0.35,yy+3.15,lw-0.7,0.85,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.1)
txt(s,lx+0.6,yy+3.28,lw-1.1,0.6,[[R("Now back to switching partners for each Epic",12.5,SAGE_DK,True,False),
    R(" — shared ownership across the stack.",12.5,INK,False,False)]],ls=1.08,anchor=MSO_ANCHOR.MIDDLE)
# communication
rx=6.9; rw=5.53
rect(s,rx,yy,rw,4.2,fill=SURFACE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.04)
rect(s,rx,yy,rw,0.7,fill=SAGE_DK,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,adj=0.06)
txt(s,rx+0.3,yy+0.06,rw-0.5,0.6,[[R("COMMUNICATION",13,WHITE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
comms=[("Daily stand-ups","Quick sync, unblock fast."),
       ("Focused silent work","Deep-work sessions, no noise."),
       ("Frequent reviews & testing","Catch issues early, together.")]
for i,(t,b) in enumerate(comms):
    cy=yy+1.05+i*1.0
    rect(s,rx+0.35,cy+0.03,0.34,0.34,fill=SAGE,shape=MSO_SHAPE.OVAL)
    txt(s,rx+0.35,cy-0.01,0.34,0.34,[[R(str(i+1),13,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,rx+0.85,cy-0.03,rw-1.2,0.85,[[R(t,14,INK,True,False)],[R(b,12,INK_SOFT,False,False)]],space_after=1,ls=1.05)
footer(s,7)
notes(s,"Both",
      "Split: Jenny owns UI and MVP refinement, Stuart owns backend and documentation — but we're back to "
      "switching partners per Epic so ownership stays shared. Rhythm: daily stand-ups, focused silent work "
      "sessions, and frequent reviews and testing.")

# ================= SLIDE 8 — Priorities
s=slide(); contentbg(s); kicker(s,0.9,0.6,"Priorities — this & next week"); heading(s,0.9,0.92,11.5,"Main goal: live testing on our phones"); rule(s,0.92,1.62)
txt(s,0.9,1.9,11.5,0.4,[[R("Everything this sprint points at one thing: ",14,INK_SOFT,False,False),
      R("getting the full flow onto real devices and demo-ready.",14,INK,True,False)]])
tiles=[("End-to-end","Stable workflow","Upload → gaps → recs, no breaks"),
       ("Onboarding","Better user flow","Smoother first-run experience"),
       ("Quality","Higher-quality recs","More relevant, better reasons"),
       ("Extra value","Recommendation ranking","Order picks by impact")]
tw=2.86;x0=0.9;yy=2.5
for i,(big,lab,sub) in enumerate(tiles):
    stat(s,x0+i*(tw+0.24),yy,tw,2.0,big,lab,sub,accent=SAGE_DK if i==0 else SAGE)
rect(s,0.9,4.85,11.53,1.35,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
txt(s,1.2,5.02,11,1.05,
    [[R("Demo readiness is the finish line.  ",13.5,SAGE_DK,True,False),
      R("We want to run the whole loop from our own phones — that's the fastest way to feel the rough edges, "
        "validate the recommendations, and get the product ready to show.",13.5,INK,False,False)]],
    ls=1.15,anchor=MSO_ANCHOR.MIDDLE)
footer(s,8)
notes(s,"Both",
      "Priority for this and next week is live testing on our own phones. Focus areas: a stable end-to-end "
      "workflow, better onboarding and user flow, higher-quality recommendations, an extra-value feature like "
      "recommendation ranking, and overall demo readiness. Getting it on real devices is how we find the "
      "rough edges fast.")

# ================= SLIDE 9 — Thank you
s=slide(); bg(s,CANVAS)
rect(s,0,0,13.333,0.14,fill=SAGE)
rect(s,0,7.36,13.333,0.14,fill=SAGE)
logo(s,0.9,0.9)
txt(s,0.9,2.5,11.5,1.2,[[R("Thank you.",50,INK,True,False)]])
rule(s,0.92,3.7,3.4)
txt(s,0.92,3.95,11,0.6,[[R("Upload your receipt. Get a smarter next shopping trip.",18,SAGE_DK,False,True)]])
txt(s,0.9,4.95,11.5,1.0,
    [[R("Jennifer Rake",15,INK,True,False),R("   &   ",15,INK_SOFT,False,False),R("Stuart Kasemeier",15,INK,True,False)],
     [R("NutriWise  ·  AI PM Capstone  ·  PM Meeting 3  —  questions & feedback welcome",12.5,INK_SOFT,False,False)]],space_after=8)
notes(s,"Both",
      "Close on the one-liner and open for questions. Headline to leave in the room: the core loop works, "
      "matching is the honest challenge we're actively de-risking, and the near-term goal is live testing on "
      "our phones toward demo readiness.")

out=os.path.join(OUTDIR,"NutriWise_PM_Meeting3.pptx")
prs.save(out); print("saved",out,"slides:",len(prs.slides._sldIdLst))
